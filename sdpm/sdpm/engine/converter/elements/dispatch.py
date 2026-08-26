# Copyright Amazon.com, Inc. or its affiliates. All Rights Reserved.
# SPDX-License-Identifier: MIT-0
"""Shape-type routing, raw passthrough, and group recursion."""
import sys

from pptx.enum.shapes import MSO_SHAPE_TYPE

from ..constants import _NS, get_emu_per_px, _base_element, _hex
from ..xml_helpers import _extract_fill_from_xml
from .shapes import extract_line_element, extract_freeform_element, extract_shape_element
from .textbox import extract_textbox_element
from .media import (extract_video_element, extract_picture_element,
                    _extract_blipfill_image, _save_referenced_images)

def _shape_needs_raw_passthrough(shape):
    """WordArt-class decoration the JSON schema can't express.

    - prstTxWarp: warped text (arch / circle / wave WordArt)
    - run-level blipFill: characters painted with a picture
    """
    try:
        x_el = shape._element
        warp = x_el.find(f'.//{{{_NS["a"]}}}prstTxWarp')
        if warp is not None and warp.get('prst') not in (None, 'textNoShape'):
            return True
        tx_body = x_el.find(f'.//{{{_NS["p"]}}}txBody')
        if tx_body is not None:
            for rpr in tx_body.iter(f'{{{_NS["a"]}}}rPr'):
                if rpr.find(f'{{{_NS["a"]}}}blipFill') is not None:
                    return True
    except Exception:
        pass
    return False


def _extract_raw_shape(shape, output_dir, slide_idx, img_counter):
    """Save shape XML verbatim (plus referenced images) for lossless rebuild."""
    from lxml import etree as _et
    elem = _base_element(shape, "rawShape")
    elem["_shapeXml"] = _et.tostring(shape._element, encoding='unicode')
    rid_map, img_counter = _save_referenced_images(shape, output_dir, slide_idx, img_counter, "raw")
    if rid_map:
        elem["_shapeImages"] = rid_map
    return elem, img_counter


def _dispatch_shape(shape, theme_colors=None, color_mapping=None, theme_styles=None, output_dir=None, slide_idx=0, img_counter=0, builder_text_color=None, pptx_path=None):
    """Dispatch shape extraction by type. Returns (elem, img_counter)."""
    elem = None
    if shape.shape_type in (MSO_SHAPE_TYPE.TEXT_BOX, MSO_SHAPE_TYPE.AUTO_SHAPE,
                            MSO_SHAPE_TYPE.FREEFORM) and _shape_needs_raw_passthrough(shape):
        return _extract_raw_shape(shape, output_dir, slide_idx, img_counter)
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        elem, img_counter = extract_group_element(shape, theme_colors, color_mapping, theme_styles, output_dir, slide_idx, img_counter, builder_text_color=builder_text_color)
    elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
        elem = _extract_blipfill_image(shape, output_dir, slide_idx, img_counter)
        if elem:
            img_counter += 1
        else:
            elem = extract_textbox_element(shape, theme_colors, color_mapping, theme_styles, builder_text_color=builder_text_color)
    elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        elem = extract_picture_element(shape, output_dir, slide_idx, img_counter, theme_colors, color_mapping)
        if elem:
            img_counter += 1
    elif shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER and shape._element.tag.endswith('}pic'):
        elem = extract_picture_element(shape, output_dir, slide_idx, img_counter, theme_colors, color_mapping)
        if elem:
            img_counter += 1
    elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        from ..table import extract_table_element
        elem = extract_table_element(shape, theme_colors, color_mapping, pptx_path)
    elif hasattr(shape, 'has_chart') and shape.has_chart:
        from ..chart import extract_chart_element
        elem = extract_chart_element(shape, theme_colors, color_mapping)
    elif shape.shape_type == MSO_SHAPE_TYPE.LINE:
        elem = extract_line_element(shape, theme_colors, color_mapping, theme_styles)
    elif shape.shape_type == 16:  # MEDIA (video)
        elem = extract_video_element(shape, output_dir, slide_idx, img_counter)
        if elem:
            img_counter += 1
    elif shape.shape_type in (MSO_SHAPE_TYPE.AUTO_SHAPE, MSO_SHAPE_TYPE.FREEFORM):
        elem = _extract_blipfill_image(shape, output_dir, slide_idx, img_counter)
        if elem:
            img_counter += 1
            return elem, img_counter
        if shape.shape_type == MSO_SHAPE_TYPE.FREEFORM:
            elem = extract_freeform_element(shape, theme_colors, color_mapping, builder_text_color=builder_text_color)
        if not elem:
            # Check if AUTO_SHAPE is actually a line (prst=line)
            try:
                prst = shape._element.spPr.find(f'.//{{{_NS["a"]}}}prstGeom')
                if prst is not None and prst.get('prst') == 'line':
                    elem = extract_line_element(shape, theme_colors, color_mapping, theme_styles)
            except Exception:
                pass
        if not elem:
            elem = extract_shape_element(shape, theme_colors, color_mapping, theme_styles, builder_text_color=builder_text_color)
    return elem, img_counter

def extract_group_element(shape, theme_colors=None, color_mapping=None, theme_styles=None, output_dir=None, slide_idx=0, img_counter=0, builder_text_color=None):
    """Extract group as element dict with nested elements.
    
    Note: python-pptx returns absolute slide coordinates for grouped shapes.
    """
    emu_per_px = get_emu_per_px()
    elem = _base_element(shape, "group", elements=[])
    # Move rotation after elements for consistent key order
    rot = elem.pop("rotation", None)
    if rot is not None:
        elem["rotation"] = rot

    # Save raw XML for groups that can't be losslessly flattened:
    # rotated groups, groups containing freeforms (recursively — vector
    # icon art nests them deep), and groups whose child coordinate space
    # is sub-pixel (px-rounded flattening collapses everything to 0x0).
    def _has_freeforms(g):
        for child in g.shapes:
            if child.shape_type == 6 and _has_freeforms(child):
                return True
            if child._element.tag.endswith('}sp') and child.shape_type == 5:
                return True
        return False

    def _subpixel_children(g):
        xf = g._element.find(f'{{{_NS["p"]}}}grpSpPr/{{{_NS["a"]}}}xfrm')
        che = xf.find(f'{{{_NS["a"]}}}chExt') if xf is not None else None
        if che is None:
            return False
        # chExt in EMU: below ~1px per unit means children live in a
        # miniature coordinate system that px rounding destroys.
        try:
            return 0 < int(che.get('cx', '0')) < int(emu_per_px * 10) or \
                   0 < int(che.get('cy', '0')) < int(emu_per_px * 10)
        except Exception:
            return False

    has_freeforms = _has_freeforms(shape)
    if rot is not None or has_freeforms or _subpixel_children(shape):
        try:
            from lxml import etree as _et
            elem["_groupXml"] = _et.tostring(shape._element, encoding='unicode')
            # Save any images referenced from inside the group XML — see
            # _save_referenced_images for why the rId mapping is required.
            rid_map, img_counter = _save_referenced_images(shape, output_dir, slide_idx, img_counter, "grp")
            if rid_map:
                elem["_groupImages"] = rid_map
        except Exception:
            pass
    
    # Extract group fill (for grpFill inheritance)
    grp_fill_color = None
    grp_fill_gradient = None
    grp_sp_pr = shape._element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}grpSpPr')
    if grp_sp_pr is not None:
        sf = grp_sp_pr.find(f'{{{_NS["a"]}}}solidFill')
        gf = grp_sp_pr.find(f'{{{_NS["a"]}}}gradFill')
        if sf is not None:
            srgb = sf.find(f'{{{_NS["a"]}}}srgbClr')
            scheme = sf.find(f'{{{_NS["a"]}}}schemeClr')
            if srgb is not None:
                grp_fill_color = _hex(srgb)
            elif scheme is not None:
                from ..color import _resolve_color_with_transforms
                grp_fill_color = _resolve_color_with_transforms(scheme, theme_colors, color_mapping)
        elif gf is not None:
            fill_info = _extract_fill_from_xml(grp_sp_pr, theme_colors, color_mapping)
            grp_fill_gradient = fill_info.get("gradient")
            # For gradient grpFill, preserve entire group XML for lossless roundtrip
            from lxml import etree as _et
            elem["_groupXml"] = _et.tostring(shape._element, encoding='unicode')
    
    # Extract each shape in the group
    for sub_shape in shape.shapes:
        try:
            sub_elem = None
            
            # Handle nested groups recursively
            sub_elem, img_counter = _dispatch_shape(sub_shape, theme_colors, color_mapping, theme_styles, output_dir, slide_idx, img_counter, builder_text_color=builder_text_color)
            
            if sub_elem:
                # Resolve grpFill: if sub-element has fill=none but XML has grpFill, use group fill
                if sub_elem.get("fill") in (None, "none"):
                    sub_sp = sub_shape._element.find(f'{{{_NS["p"]}}}spPr')
                    if sub_sp is not None and sub_sp.find(f'{{{_NS["a"]}}}grpFill') is not None:
                        if grp_fill_color:
                            sub_elem["fill"] = grp_fill_color
                        elif grp_fill_gradient:
                            sub_elem["gradient"] = grp_fill_gradient
                # Propagate grpFill to nested group children
                if sub_elem.get("type") == "group":
                    grp_sp = sub_shape._element.find(f'{{{_NS["p"]}}}grpSpPr')
                    if grp_sp is not None and grp_sp.find(f'{{{_NS["a"]}}}grpFill') is not None:
                        for child_el in sub_elem.get("elements", []):
                            if child_el.get("fill") in (None, "none"):
                                if grp_fill_color:
                                    child_el["fill"] = grp_fill_color
                                elif grp_fill_gradient:
                                    child_el["gradient"] = grp_fill_gradient

                # Transform coordinates from child coordinate system to slide coordinates
                grp_sp_pr = shape._element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}grpSpPr')
                xfrm = grp_sp_pr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm') if grp_sp_pr is not None else None
                
                if xfrm is not None:
                    off = xfrm.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}off')
                    ext = xfrm.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}ext')
                    ch_off = xfrm.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}chOff')
                    ch_ext = xfrm.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}chExt')
                    
                    if off is not None and ch_off is not None and ext is not None and ch_ext is not None:
                        group_off_x = int(off.get('x'))
                        group_off_y = int(off.get('y'))
                        group_ext_cx = int(ext.get('cx'))
                        group_ext_cy = int(ext.get('cy'))
                        ch_off_x = int(ch_off.get('x'))
                        ch_off_y = int(ch_off.get('y'))
                        ch_ext_cx = int(ch_ext.get('cx'))
                        ch_ext_cy = int(ch_ext.get('cy'))
                        
                        # Transform: abs = group_off + (child - chOff) * (group_ext / ch_ext)
                        child_x = sub_shape.left
                        child_y = sub_shape.top
                        
                        scale_x = group_ext_cx / ch_ext_cx if ch_ext_cx != 0 else 1
                        scale_y = group_ext_cy / ch_ext_cy if ch_ext_cy != 0 else 1
                        
                        abs_x = group_off_x + (child_x - ch_off_x) * scale_x
                        abs_y = group_off_y + (child_y - ch_off_y) * scale_y
                        
                        if sub_elem.get("type") == "line" and "x1" in sub_elem:
                            # Lines carry endpoints, not x/y/width/height —
                            # transform x1/y1/x2/y2 through the group xfrm.
                            def _gx(px_):
                                return round((group_off_x + (px_ * emu_per_px - ch_off_x) * scale_x) / emu_per_px)
                            def _gy(py_):
                                return round((group_off_y + (py_ * emu_per_px - ch_off_y) * scale_y) / emu_per_px)
                            for k in ("x1", "x2"):
                                if sub_elem.get(k) is not None:
                                    sub_elem[k] = _gx(sub_elem[k])
                            for k in ("y1", "y2"):
                                if sub_elem.get(k) is not None:
                                    sub_elem[k] = _gy(sub_elem[k])
                            elem["elements"].append(sub_elem)
                            continue

                        sub_elem["x"] = round(abs_x / emu_per_px)
                        sub_elem["y"] = round(abs_y / emu_per_px)
                        sub_elem["width"] = round(sub_shape.width * scale_x / emu_per_px)
                        sub_elem["height"] = round(sub_shape.height * scale_y / emu_per_px)
                        # For freeform in group: drop raw path XML, let builder reconstruct
                        # from px coords (which match the group-scaled shape size)
                        if sub_elem.get("type") == "freeform":
                            sub_elem["_xEmu"] = round(abs_x)
                            sub_elem["_yEmu"] = round(abs_y)
                            sub_elem["_widthEmu"] = round(sub_shape.width * scale_x)
                            sub_elem["_heightEmu"] = round(sub_shape.height * scale_y)
                            sub_elem.pop("_pathLstXml", None)
                        
                        # For nested groups, also transform all children recursively
                        if sub_elem.get("type") == "group" and sub_elem.get("elements"):
                            def _apply_group_transform(elements, gox, goy, gcx, gcy, sx, sy):
                                for el in elements:
                                    if el.get("type") == "line" and "x1" in el:
                                        for k in ("x1", "x2"):
                                            if el.get(k) is not None:
                                                el[k] = round((gox + (el[k] * emu_per_px - gcx) * sx) / emu_per_px)
                                        for k in ("y1", "y2"):
                                            if el.get(k) is not None:
                                                el[k] = round((goy + (el[k] * emu_per_px - gcy) * sy) / emu_per_px)
                                        continue
                                    if "x" in el and "y" in el:
                                        old_x = el["x"] * emu_per_px
                                        old_y = el["y"] * emu_per_px
                                        new_x = gox + (old_x - gcx) * sx
                                        new_y = goy + (old_y - gcy) * sy
                                        el["x"] = round(new_x / emu_per_px)
                                        el["y"] = round(new_y / emu_per_px)
                                        if el.get("type") == "freeform":
                                            el["_xEmu"] = round(new_x)
                                            el["_yEmu"] = round(new_y)
                                    if "width" in el:
                                        el["width"] = round(el["width"] * sx)
                                    if "height" in el:
                                        el["height"] = round(el["height"] * sy)
                                    if el.get("type") == "freeform":
                                        if el.get("_widthEmu"):
                                            el["_widthEmu"] = round(el["_widthEmu"] * sx)
                                        else:
                                            el["_widthEmu"] = round(el["width"] * emu_per_px)
                                        if el.get("_heightEmu"):
                                            el["_heightEmu"] = round(el["_heightEmu"] * sy)
                                        else:
                                            el["_heightEmu"] = round(el["height"] * emu_per_px)
                                        el.pop("_pathLstXml", None)
                                    if el.get("type") == "group" and el.get("elements"):
                                        _apply_group_transform(el["elements"], gox, goy, gcx, gcy, sx, sy)
                            _apply_group_transform(sub_elem["elements"], group_off_x, group_off_y, ch_off_x, ch_off_y, scale_x, scale_y)
                    else:
                        # Fallback: use python-pptx coordinates as-is
                        pass
                
                # Propagate group rotation to child elements
                if shape.rotation != 0:
                    child_rot = sub_elem.get("rotation", 0)
                    sub_elem["rotation"] = round(child_rot + shape.rotation, 1)
                
                elem["elements"].append(sub_elem)
        except Exception as e:
            print(f"Warning: Failed to extract grouped shape: {e}", file=sys.stderr)
    
    return elem, img_counter
