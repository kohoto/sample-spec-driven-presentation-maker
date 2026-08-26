# Copyright Amazon.com, Inc. or its affiliates. All Rights Reserved.
# SPDX-License-Identifier: MIT-0
"""Textbox extraction (delegates preset-geometry text boxes to shapes)."""
import json

from ..constants import _NS, get_emu_per_px, _serialize_lstStyle, _hex, _add_flip
from ..xml_helpers import (_extract_fill_from_xml, _extract_line_from_xml,
                           _extract_effects_from_xml, _extract_visual_effects)
from ..text import _extract_styled_text, _detect_font_size, _get_alignment
from .shapes import extract_shape_element

def extract_textbox_element(shape, theme_colors=None, color_mapping=None, theme_styles=None, is_placeholder=False, builder_text_color=None):
    """Extract textbox as element dict."""
    emu_per_px = get_emu_per_px()
    # Check if it's actually a shape with preset geometry (not a plain textbox)
    try:
        sp_pr = shape._element.spPr
        prst_geom = sp_pr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom')
        if prst_geom is not None:
            prst = prst_geom.get('prst')
            # If it has any preset geometry (not just 'rect'), treat as shape
            if prst and prst != 'rect':
                # This is a shape with text, not a plain textbox
                return extract_shape_element(shape, theme_colors, color_mapping, theme_styles, builder_text_color=builder_text_color)
    except Exception:
        pass
    
    elem = {
        "type": "textbox",
        "x": round(shape.left / emu_per_px),  # px (1920x1080 basis)
        "y": round(shape.top / emu_per_px),
        "width": round(shape.width / emu_per_px),
    }
    
    # Extract height (for TEXT_TO_FIT_SHAPE auto-shrink)
    if shape.height:
        h_px = round(shape.height / emu_per_px)
        if h_px > 10:
            elem["height"] = h_px
    # Extract rotation
    if shape.rotation != 0:
        elem["rotation"] = round(shape.rotation, 1)
    
    # Extract flip
    _add_flip(elem, shape)
    
    # Extract autoWidth
    try:
        body_pr = shape._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr')
        if body_pr is not None:
            if body_pr.get('wrap') == 'none':
                elem["autoWidth"] = True
            vert = body_pr.get('vert')
            if vert:
                elem["textDirection"] = vert
    except Exception:
        pass
    
    # Extract margins (EMU → px)
    # Builder default for textbox: left/right=0, top/bottom=PowerPoint default
    tf = shape.text_frame
    if tf.margin_left is not None and tf.margin_left != 0:
        elem["marginLeft"] = round(tf.margin_left / emu_per_px)
    if tf.margin_top is not None and tf.margin_top != 45720:
        elem["marginTop"] = round(tf.margin_top / emu_per_px)
    if tf.margin_right is not None and tf.margin_right != 0:
        elem["marginRight"] = round(tf.margin_right / emu_per_px)
    if tf.margin_bottom is not None and tf.margin_bottom != 45720:
        elem["marginBottom"] = round(tf.margin_bottom / emu_per_px)
    
    # Extract vertical anchor (builder textbox default is top when unset)
    if tf.vertical_anchor is not None:
        _va_reverse = {1: "top", 3: "middle", 4: "bottom"}
        va = _va_reverse.get(int(tf.vertical_anchor))
        if va:
            elem["verticalAlign"] = va

    # Extract fill and line using XML helpers
    try:
        sp_pr_xml = shape._element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}spPr')
        elem.update(_extract_fill_from_xml(sp_pr_xml, theme_colors, color_mapping))
        elem.update(_extract_line_from_xml(sp_pr_xml, theme_colors, color_mapping))
        elem.update(_extract_visual_effects(sp_pr_xml, theme_colors, color_mapping))
    except Exception:
        pass
    
    # Extract textGradient from runs with gradFill
    try:
        grad_runs = []
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                rpr = run._r.find(f'{{{_NS["a"]}}}rPr')
                if rpr is not None:
                    grad = rpr.find(f'{{{_NS["a"]}}}gradFill')
                    if grad is not None:
                        stops = []
                        for gs in grad.findall(f'.//{{{_NS["a"]}}}gs'):
                            pos = round(int(gs.get('pos', '0')) / 100000, 2)
                            srgb = gs.find(f'{{{_NS["a"]}}}srgbClr')
                            if srgb is not None:
                                stops.append({"position": pos, "color": _hex(srgb)})
                        if stops:
                            angle = 0
                            lin = grad.find(f'{{{_NS["a"]}}}lin')
                            if lin is not None:
                                angle = round(int(lin.get('ang', '0')) / 60000)
                            grad_runs.append({"text": run.text, "gradient": {"angle": angle, "stops": stops}})
        if grad_runs:
            # Count total runs with text
            total_runs = sum(1 for p in shape.text_frame.paragraphs for r in p.runs if r.text.strip())
            grads = [json.dumps(gr["gradient"], sort_keys=True) for gr in grad_runs]
            # Promote to textGradient only if ALL runs have the same gradient
            if len(set(grads)) == 1 and len(grad_runs) >= total_runs:
                elem["textGradient"] = grad_runs[0]["gradient"]
            else:
                elem["_textGradientRuns"] = grad_runs
    except Exception:
        pass

    # Extract run-level text effects (glow/shadow on the characters). All
    # runs sharing one effectLst is the common case (decorated headline);
    # store the raw XML for lossless rebuild.
    try:
        from lxml import etree as _et_eff
        effect_xmls = set()
        has_run = False
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if not run.text.strip():
                    continue
                has_run = True
                rpr = run._r.find(f'{{{_NS["a"]}}}rPr')
                eff = rpr.find(f'{{{_NS["a"]}}}effectLst') if rpr is not None else None
                if eff is not None and len(eff) > 0:
                    effect_xmls.add(_et_eff.tostring(eff, encoding='unicode'))
                else:
                    effect_xmls.add("")
        if has_run and len(effect_xmls) == 1:
            xml = effect_xmls.pop()
            if xml:
                elem["_textEffects"] = xml
    except Exception:
        pass

    # Detect cap=none and bold=off overrides (when lstStyle has cap=all / b=1)
    try:
        _all_runs = [r for p in shape.text_frame.paragraphs for r in p.runs]
        if _all_runs:
            if all(r._r.find(f'{{{_NS["a"]}}}rPr') is not None and
                   r._r.find(f'{{{_NS["a"]}}}rPr').get('cap') == 'none'
                   for r in _all_runs):
                elem["_capNone"] = True
            if all(r._r.find(f'{{{_NS["a"]}}}rPr') is not None and
                   r._r.find(f'{{{_NS["a"]}}}rPr').get('b') == '0'
                   for r in _all_runs):
                elem["_boldOff"] = True
    except Exception:
        pass

    # Extract text with styles
    text_parts = []
    default_font_size = None
    
    # Determine default text color (must match builder's theme_colors["text"])
    # For placeholders, don't set default_text_color — lstStyle defines the actual default
    default_text_color = None
    if not is_placeholder:
        default_text_color = builder_text_color
        if not default_text_color and color_mapping and theme_colors:
            tx1_mapped = color_mapping.get('tx1', 'dk1')
            default_text_color = theme_colors.get(tx1_mapped)
    
    # Check if multiple paragraphs (should be items array)
    paragraphs_with_text = [p for p in shape.text_frame.paragraphs if p.text.strip()]
    all_paragraphs = list(shape.text_frame.paragraphs)
    has_lstStyle = _serialize_lstStyle(shape) is not None
    
    if len(all_paragraphs) > 1:
        # Multiple paragraphs - extract as paragraphs with bullet info
        default_font_size = None if (is_placeholder or has_lstStyle) else _detect_font_size(all_paragraphs)
        paragraphs = []
        for paragraph in all_paragraphs:
            
            # Empty paragraph
            if not paragraph.text.strip():
                paragraphs.append({"text": ""})
                continue
            
            # Check for bullet or numbering
            has_bullet = False
            numbering_type = None
            bu_font = None
            mar_l = None
            indent = None
            space_after = None
            space_before = None
            line_spacing = None
            try:
                pPr = paragraph._element.pPr
                if pPr is not None:
                    bu_auto_num = pPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}buAutoNum')
                    bu_char = pPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}buChar')
                    
                    if bu_auto_num is not None:
                        numbering_type = bu_auto_num.get('type', 'arabicPeriod')
                    elif bu_char is not None:
                        has_bullet = True
                    
                    bu_font_elem = pPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}buFont')
                    if bu_font_elem is not None:
                        bu_font = bu_font_elem.get('typeface')
                    mar_l = pPr.get('marL')
                    indent = pPr.get('indent')
                    spc_aft = pPr.find('.//a:spcAft/a:spcPts', _NS)
                    if spc_aft is not None:
                        space_after = spc_aft.get('val')
                    spc_bef = pPr.find('.//a:spcBef/a:spcPts', _NS)
                    if spc_bef is not None:
                        space_before = spc_bef.get('val')
                    ln_spc = pPr.find('.//a:lnSpc/a:spcPts', _NS)
                    if ln_spc is not None:
                        line_spacing = ('pts', ln_spc.get('val'))
                    else:
                        ln_spc_pct = pPr.find('.//a:lnSpc/a:spcPct', _NS)
                        if ln_spc_pct is not None:
                            line_spacing = ('pct', ln_spc_pct.get('val'))
            except Exception:
                pass
            
            item_text = _extract_styled_text(paragraph.runs, theme_colors, color_mapping, default_font_size=default_font_size, default_text_color=default_text_color, is_placeholder=is_placeholder, paragraph=paragraph)
            para_info = {"text": item_text}
            # Explicit paragraph alignment — without it a shape-level
            # lstStyle default (e.g. centered) silently wins.
            _algn = _get_alignment(paragraph)
            if _algn:
                para_info["align"] = _algn
            if has_bullet or numbering_type:
                list_def = {}
                if numbering_type:
                    list_def["type"] = numbering_type
                else:
                    list_def["type"] = "disc"
                level = paragraph.level if paragraph.level else 0
                if level > 0:
                    list_def["level"] = level
                para_info["list"] = list_def
            if bu_font:
                para_info["buFont"] = bu_font
            if mar_l is not None:
                para_info["marL"] = int(mar_l)
            if indent is not None:
                para_info["indent"] = int(indent)
            if space_after is not None:
                para_info["spaceAfter"] = int(space_after)
            if space_before is not None:
                para_info["spaceBefore"] = int(space_before)
            if line_spacing:
                if line_spacing[0] == 'pct':
                    para_info["lineSpacingPct"] = int(line_spacing[1])
                else:
                    para_info["lineSpacing"] = int(line_spacing[1])
            
            # Paragraph level (for sub-bullets)
            try:
                pPr = paragraph._element.pPr
                if pPr is not None:
                    lvl = pPr.get('lvl')
                    if lvl and lvl != '0':
                        para_info["level"] = int(lvl)
            except Exception:
                pass
            
            paragraphs.append(para_info)
        
        if paragraphs:
            elem["paragraphs"] = paragraphs
            
            # Add fontSize if not default
            if default_font_size and default_font_size != 18:
                elem["fontSize"] = default_font_size
            
            # Get alignment - per paragraph if mixed, top-level if uniform
            aligns = [_get_alignment(p) for p in paragraphs_with_text]
            unique = set(a for a in aligns if a)
            if len(unique) <= 1:
                align = aligns[0] if aligns else None
                if align and align != "left":
                    elem["align"] = align
            else:
                # Mixed alignment: set per-paragraph
                for para_info, paragraph in zip(paragraphs, shape.text_frame.paragraphs):
                    a = _get_alignment(paragraph)
                    if a:
                        para_info["align"] = a
            
            # Preserve lstStyle for roundtrip fidelity
            _lst = _serialize_lstStyle(shape) if shape.has_text_frame else None
            if _lst:
                elem["_lstStyle"] = _lst
            
            # Extract character spacing
            _spc_vals = set()
            for _p in shape.text_frame.paragraphs:
                for _r in _p.runs:
                    _rPr = _r._r.find('{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
                    _s = _rPr.get('spc') if _rPr is not None else None
                    if _s:
                        _spc_vals.add(int(_s))
            if len(_spc_vals) == 1:
                elem["_spc"] = _spc_vals.pop()
            
            return elem
    
    # Single paragraph - extract as text
    default_font_size = None if (is_placeholder or has_lstStyle) else _detect_font_size(shape.text_frame.paragraphs)
    for paragraph in shape.text_frame.paragraphs:
        text_parts.append(_extract_styled_text(paragraph.runs, theme_colors, color_mapping, default_font_size=default_font_size, default_text_color=default_text_color, is_placeholder=is_placeholder, paragraph=paragraph))
    
    elem["text"] = ''.join(text_parts)

    # endParaRPr pins the paragraph line height (e.g. a full-size 80pt
    # endParaRPr next to a baseline-shrunk run keeps the line tall;
    # dropping it shifts the text up within the box).
    if shape.text_frame.paragraphs:
        _last_p = shape.text_frame.paragraphs[-1]
        _endPr = _last_p._element.find(f'{{{_NS["a"]}}}endParaRPr')
        if _endPr is not None and _endPr.get('sz'):
            _end_sz = int(_endPr.get('sz')) / 100
            _last_runs = _last_p.runs
            _last_run_sz = (_last_runs[-1].font.size.pt
                            if _last_runs and _last_runs[-1].font.size else None)
            _has_baseline = any(
                (r._r.find(f'{{{_NS["a"]}}}rPr') is not None
                 and r._r.find(f'{{{_NS["a"]}}}rPr').get('baseline'))
                for r in _last_runs)
            if _has_baseline or (_last_run_sz is not None and _end_sz != _last_run_sz):
                elem["_endParaSize"] = _end_sz

    # Extract indent/marL from first paragraph
    if shape.text_frame.paragraphs:
        from pptx.oxml.ns import qn as _qn
        pPr = shape.text_frame.paragraphs[0]._element.find(_qn('a:pPr'))
        if pPr is not None:
            _indent = pPr.get('indent')
            if _indent is not None:
                elem["indent"] = int(_indent)
            _marL = pPr.get('marL')
            if _marL is not None:
                elem["marL"] = int(_marL)
    
    # Add fontSize if consistent
    if default_font_size:
        elem["fontSize"] = default_font_size
    
    # Detect alignment
    if shape.text_frame.paragraphs:
        align = _get_alignment(shape.text_frame.paragraphs[0])
        if align:
            elem["align"] = align
        # Line spacing from first paragraph
        pPr = shape.text_frame.paragraphs[0]._element.find('{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
        if pPr is not None:
            lnSpc_pct = pPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}lnSpc/{http://schemas.openxmlformats.org/drawingml/2006/main}spcPct')
            if lnSpc_pct is not None:
                elem["lineSpacingPct"] = int(lnSpc_pct.get('val'))
            # Fixed-point spacing (spcPts) — e.g. a 48pt title with 31.2pt
            # spacing renders much higher/tighter than the default.
            lnSpc_pts = pPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}lnSpc/{http://schemas.openxmlformats.org/drawingml/2006/main}spcPts')
            if lnSpc_pts is not None:
                elem["lineSpacingPt"] = int(lnSpc_pts.get('val')) / 100
    
    # Extract visual effects
    try:
        sp_pr_xml = shape._element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}spPr')
        if sp_pr_xml is None:
            sp_pr_xml = shape._element.spPr if hasattr(shape._element, 'spPr') else None
        elem.update(_extract_effects_from_xml(sp_pr_xml, theme_colors, color_mapping))
    except Exception:
        pass
    
    # Preserve lstStyle for roundtrip fidelity
    _lst = _serialize_lstStyle(shape) if shape.has_text_frame else None
    if _lst:
        elem["_lstStyle"] = _lst
    
    # Extract character spacing (spc) if uniform across all runs
    if shape.has_text_frame:
        spc_values = set()
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                rPr = r._r.find('{http://schemas.openxmlformats.org/drawingml/2006/main}rPr')
                spc = rPr.get('spc') if rPr is not None else None
                if spc:
                    spc_values.add(int(spc))
        if len(spc_values) == 1:
            elem["_spc"] = spc_values.pop()
    
    return elem

