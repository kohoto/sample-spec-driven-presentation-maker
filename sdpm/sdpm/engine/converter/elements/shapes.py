# Copyright Amazon.com, Inc. or its affiliates. All Rights Reserved.
# SPDX-License-Identifier: MIT-0
"""Shape-name resolution and line / freeform / auto-shape extraction."""
import sys

from pptx.enum.shapes import MSO_SHAPE

from ..constants import _NS, get_emu_per_px, _base_element, _add_flip, _serialize_lstStyle
from ..color import _resolve_color_with_transforms
from ..xml_helpers import (extract_line_dash, _resolve_line_from_style, _extract_fill_from_xml,
                           _extract_line_from_xml, _extract_visual_effects)
from ..text import _extract_shape_text

_SHAPE_MAP = {
    MSO_SHAPE.RECTANGLE: "rectangle", MSO_SHAPE.ROUNDED_RECTANGLE: "rounded_rectangle",
    MSO_SHAPE.OVAL: "oval", MSO_SHAPE.RIGHT_ARROW: "arrow_right",
    MSO_SHAPE.LEFT_ARROW: "arrow_left", MSO_SHAPE.UP_ARROW: "arrow_up",
    MSO_SHAPE.DOWN_ARROW: "arrow_down", MSO_SHAPE.ISOSCELES_TRIANGLE: "triangle",
    MSO_SHAPE.DIAMOND: "diamond", MSO_SHAPE.PENTAGON: "pentagon",
    MSO_SHAPE.HEXAGON: "hexagon", MSO_SHAPE.CHEVRON: "chevron",
    MSO_SHAPE.RIGHT_BRACE: "right_brace", MSO_SHAPE.LEFT_BRACE: "left_brace",
    60: "arrow_circular",
}
_PRESET_MAP = {
    'roundRect': 'rounded_rectangle', 'rect': 'rectangle', 'ellipse': 'oval',
    'triangle': 'triangle', 'diamond': 'diamond', 'pentagon': 'pentagon',
    'hexagon': 'hexagon', 'chevron': 'chevron', 'homePlate': 'pentagon',
    'heart': 'heart', 'cloud': 'cloud', 'lightningBolt': 'lightning_bolt',
    'star5': 'star_5_point', 'noSmoking': 'no_symbol', 'cross': 'cross', 'plus': 'cross',
    'trapezoid': 'trapezoid', 'parallelogram': 'parallelogram',
    'donut': 'donut', 'arc': 'arc', 'blockArc': 'block_arc', 'chord': 'chord',
    'pie': 'pie', 'pieWedge': 'pie_wedge',
    'leftRightArrow': 'arrow_left_right', 'upDownArrow': 'arrow_up_down',
    'curvedRightArrow': 'arrow_curved_right', 'curvedLeftArrow': 'arrow_curved_left',
    'curvedUpArrow': 'arrow_curved_up', 'curvedDownArrow': 'arrow_curved_down',
    'circularArrow': 'arrow_circular', 'leftCircularArrow': 'arrow_circular_left',
    'leftRightCircularArrow': 'arrow_circular_left_right',
    'calloutRoundRect': 'callout_rounded_rectangle', 'wedgeRoundRectCallout': 'callout_rounded_rectangle',
    'calloutRect': 'callout_rectangle', 'wedgeRectCallout': 'callout_rectangle',
    'calloutEllipse': 'callout_oval', 'wedgeEllipseCallout': 'callout_oval',
    'flowChartProcess': 'flowchart_process', 'flowChartDecision': 'flowchart_decision',
    'flowChartTerminator': 'flowchart_terminator',
    'leftBracket': 'left_bracket', 'rightBracket': 'right_bracket',
    'can': 'cylinder', 'mathNotEqual': 'math_not_equal',
}

def _resolve_shape_name(shape):
    """Resolve shape preset name from python-pptx or XML."""
    if shape.shape_type == 5:  # MSO_SHAPE_TYPE.FREEFORM
        return "rounded_rectangle"
    name = None
    try:
        if hasattr(shape, 'auto_shape_type') and shape.auto_shape_type in _SHAPE_MAP:
            name = _SHAPE_MAP[shape.auto_shape_type]
    except Exception:
        pass
    if not name:
        try:
            prst = shape._element.spPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom')
            if prst is not None:
                prst_val = prst.get('prst')
                name = _PRESET_MAP.get(prst_val, prst_val)  # Use raw prst value as fallback
        except Exception:
            pass
    if name == "oval" and shape.width == shape.height:
        return "circle"
    return name or "rounded_rectangle"

def extract_line_element(shape, theme_colors=None, color_mapping=None, theme_styles=None):
    """Extract line/connector as element dict."""
    emu_per_px = get_emu_per_px()
    try:
        # Build x1/y1/x2/y2 from bounding box + flip
        x = round(shape.left / emu_per_px)
        y = round(shape.top / emu_per_px)
        w = round(shape.width / emu_per_px)
        h = round(shape.height / emu_per_px)
        x1, y1, x2, y2 = x, y, x + w, y + h

        # Absorb flip and rotation into coordinates.
        # OOXML renders a connector inside its bounding box (start at one
        # corner, end at the opposite), flips it, then rotates the whole box
        # about its center. The schema has no rotation on lines, so bake the
        # rotation into the endpoints instead.
        rot_deg = 0
        try:
            xfrm = shape._element.spPr.find(
                './/{http://schemas.openxmlformats.org/drawingml/2006/main}xfrm')
            if xfrm is not None:
                if xfrm.get('flipH') == '1':
                    x1, x2 = x2, x1
                if xfrm.get('flipV') == '1':
                    y1, y2 = y2, y1
                rot_deg = int(xfrm.get('rot', '0')) / 60000
        except Exception:
            pass
        if rot_deg:
            import math
            theta = math.radians(rot_deg)  # clockwise in y-down coords
            c, s = math.cos(theta), math.sin(theta)
            cx0, cy0 = x + w / 2, y + h / 2
            def _rot(px_, py_):
                dx, dy = px_ - cx0, py_ - cy0
                return round(cx0 + dx * c - dy * s), round(cy0 + dx * s + dy * c)
            x1, y1 = _rot(x1, y1)
            x2, y2 = _rot(x2, y2)

        elem = {"type": "line", "x1": x1, "y1": y1, "x2": x2, "y2": y2}
        
        # Extract connector type from XML
        try:
            sp_pr = shape._element.spPr
            prst_geom = sp_pr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}prstGeom')
            if prst_geom is not None:
                prst = prst_geom.get('prst')
                if prst:
                    # Save exact preset type
                    elem["preset"] = prst
                    
                    # Map to general connector type
                    if 'straight' in prst.lower():
                        elem["connectorType"] = "straight"
                    elif 'bent' in prst.lower():
                        elem["connectorType"] = "elbow"
                        # A 90/270° rotated bent connector renders V-H-V
                        # (first segment vertical); the builder reconstructs
                        # elbows as H-V-H unless told otherwise.
                        r = rot_deg % 360
                        if 45 <= r < 135 or 225 <= r < 315:
                            elem["elbowStart"] = "vertical"
                    elif 'curved' in prst.lower():
                        elem["connectorType"] = "curved"
                
                # Extract adjustments
                av_lst = prst_geom.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}avLst')
                if av_lst is not None:
                    adjustments = []
                    for gd in av_lst.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}gd'):
                        fmla = gd.get('fmla', '')
                        if fmla.startswith('val '):
                            adj_val = int(fmla.split()[1])
                            adjustments.append(adj_val / 100000.0)
                    if adjustments:
                        elem["adjustments"] = adjustments
        except Exception:
            elem["connectorType"] = "straight"  # default
        
        # Extract arrow heads from XML
        try:
            ln = shape._element.spPr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}ln')
            if ln is not None:
                head_end = ln.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}headEnd')
                if head_end is not None:
                    head_type = head_end.get('type')
                    if head_type:
                        elem["arrowStart"] = head_type
                        if head_end.get('w'):
                            elem["arrowStartWidth"] = head_end.get('w')
                        if head_end.get('len'):
                            elem["arrowStartLength"] = head_end.get('len')

                tail_end = ln.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}tailEnd')
                if tail_end is not None:
                    tail_type = tail_end.get('type')
                    if tail_type:
                        elem["arrowEnd"] = tail_type
                        if tail_end.get('w'):
                            elem["arrowEndWidth"] = tail_end.get('w')
                        if tail_end.get('len'):
                            elem["arrowEndLength"] = tail_end.get('len')
        except Exception:
            pass
        
        # Extract line color or gradient (use XML helper)
        sp_pr_xml = shape._element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}spPr')
        line_info = _extract_line_from_xml(sp_pr_xml, theme_colors, color_mapping)
        if "line" in line_info and line_info["line"] != "none":
            elem["color"] = line_info["line"]
        elif line_info.get("line") == "none":
            elem["color"] = "none"
        else:
            # Resolve from style reference
            style_info = _resolve_line_from_style(shape, theme_colors, color_mapping, theme_styles)
            if style_info.get("line"):
                elem["color"] = style_info["line"]
            if style_info.get("lineWidth"):
                elem["lineWidth"] = style_info["lineWidth"]
        if "lineGradient" in line_info:
            elem["lineGradient"] = line_info["lineGradient"]
        if "lineWidth" in line_info:
            elem["lineWidth"] = line_info["lineWidth"]
        
        # Extract dash style
        dash = extract_line_dash(shape)
        if dash:
            elem["dashStyle"] = dash

        # No effects in source → say so explicitly (same rule as shapes).
        # python-pptx's add_connector default <p:style> has effectRef idx=1
        # (theme shadow), which painted a shadow under plain lines.
        try:
            style_el = shape._element.find(f'{{{_NS["p"]}}}style')
            eff_ref = style_el.find(f'{{{_NS["a"]}}}effectRef') if style_el is not None else None
            has_own_effects = False
            sp_pr_el = shape._element.find(f'{{{_NS["p"]}}}spPr')
            if sp_pr_el is not None:
                eff_lst = sp_pr_el.find(f'{{{_NS["a"]}}}effectLst')
                has_own_effects = eff_lst is not None and len(eff_lst) > 0
            if not has_own_effects and (
                    eff_ref is None or int(eff_ref.get('idx', '0') or 0) == 0):
                elem["_noEffects"] = True
        except Exception:
            pass

        return elem
    except Exception as e:
        print(f"Warning: Failed to extract line: {e}", file=sys.stderr)
        return None

def extract_freeform_element(shape, theme_colors=None, color_mapping=None, builder_text_color=None):
    """Extract freeform/curve shape as element dict with path commands in px."""
    emu_per_px = get_emu_per_px()
    try:
        sp_pr = shape._element.spPr
        cust_geom = sp_pr.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}custGeom')
        if cust_geom is None:
            return None

        elem = _base_element(shape, "freeform")
        _add_flip(elem, shape)
        # Preserve exact EMU size for freeform roundtrip fidelity
        elem["_widthEmu"] = shape.width
        elem["_heightEmu"] = shape.height

        # Extract all paths
        path_elements = cust_geom.findall('.//a:pathLst/a:path', _NS)
        if not path_elements:
            return None

        def _extract_path_commands(path_el):
            """Extract commands from a single path element."""
            path_w = int(path_el.get('w', shape.width))
            path_h = int(path_el.get('h', shape.height))
            sx = shape.width / path_w if path_w else 1
            sy = shape.height / path_h if path_h else 1

            def to_px(x, y):
                return round(int(x) * sx / emu_per_px, 1), round(int(y) * sy / emu_per_px, 1)

            path = []
            for child in path_el:
                tag = child.tag.split('}')[-1]
                pts = child.findall('a:pt', _NS)
                if tag == 'moveTo' and pts:
                    px, py = to_px(pts[0].get('x'), pts[0].get('y'))
                    path.append({"cmd": "M", "x": px, "y": py})
                elif tag == 'lnTo' and pts:
                    px, py = to_px(pts[0].get('x'), pts[0].get('y'))
                    path.append({"cmd": "L", "x": px, "y": py})
                elif tag == 'cubicBezTo' and len(pts) == 3:
                    coords = [to_px(p.get('x'), p.get('y')) for p in pts]
                    path.append({"cmd": "C", "pts": [[c[0], c[1]] for c in coords]})
                elif tag == 'quadBezTo' and len(pts) == 2:
                    coords = [to_px(p.get('x'), p.get('y')) for p in pts]
                    path.append({"cmd": "Q", "pts": [[c[0], c[1]] for c in coords]})
                elif tag == 'arcTo':
                    wR = int(child.get('wR', 0))
                    hR = int(child.get('hR', 0))
                    stAng = int(child.get('stAng', 0))
                    swAng = int(child.get('swAng', 0))
                    path.append({
                        "cmd": "A",
                        "wR": round(wR * sx / emu_per_px, 1),
                        "hR": round(hR * sy / emu_per_px, 1),
                        "stAng": round(stAng / 60000, 2),
                        "swAng": round(swAng / 60000, 2),
                    })
                elif tag == 'close':
                    path.append({"cmd": "Z"})
            return path

        if len(path_elements) == 1:
            # Single path → "path" key (backward compatible)
            elem["path"] = _extract_path_commands(path_elements[0])
            fill_attr = path_elements[0].get('fill')
            if fill_attr and fill_attr != 'norm':
                elem["pathFill"] = fill_attr
        else:
            # Multiple paths → "paths" key
            paths = []
            for pe in path_elements:
                p = {"commands": _extract_path_commands(pe)}
                fill_attr = pe.get('fill')
                if fill_attr and fill_attr != 'norm':
                    p["fill"] = fill_attr
                paths.append(p)
            elem["paths"] = paths

        # Preserve raw pathLst XML for lossless roundtrip
        path_el_first = path_elements[0]
        path_w = int(path_el_first.get('w', shape.width))
        path_h = int(path_el_first.get('h', shape.height))
        if path_w == shape.width and path_h == shape.height:
            from lxml import etree as _et
            pathLst = cust_geom.find('.//a:pathLst', _NS)
            if pathLst is not None:
                elem["_pathLstXml"] = _et.tostring(pathLst, encoding='unicode')

        # Fill
        elem.update(_extract_fill_from_xml(sp_pr, theme_colors, color_mapping))

        # Line
        line_info = _extract_line_from_xml(sp_pr, theme_colors, color_mapping)
        elem.update(line_info)

        # Line opacity
        ln = sp_pr.find(f'.//{{{_NS["a"]}}}ln')
        if ln is not None:
            solid = ln.find(f'{{{_NS["a"]}}}solidFill')
            if solid is not None:
                for clr_tag in ('srgbClr', 'schemeClr'):
                    clr = solid.find(f'{{{_NS["a"]}}}{clr_tag}')
                    if clr is not None:
                        alpha = clr.find(f'{{{_NS["a"]}}}alpha')
                        if alpha is not None:
                            elem["lineOpacity"] = round(int(alpha.get('val')) / 100000, 2)
                        break

        # Arrow heads
        try:
            if ln is not None:
                for attr, tag in [("headEnd", "headEnd"), ("tailEnd", "tailEnd")]:
                    el = ln.find(f'{{{_NS["a"]}}}{tag}')
                    if el is not None and el.get('type'):
                        elem[attr] = el.get('type')
        except Exception:
            pass

        # Effects
        elem.update(_extract_visual_effects(sp_pr, theme_colors, color_mapping))

        # Text (if freeform contains text)
        if shape.has_text_frame and shape.text_frame.text.strip():
            _extract_shape_text(shape, elem, theme_colors, color_mapping, builder_text_color=builder_text_color)

        return elem
    except Exception as e:
        print(f"Warning: Failed to extract freeform: {e}", file=sys.stderr)
        return None

def extract_shape_element(shape, theme_colors=None, color_mapping=None, theme_styles=None, builder_text_color=None):
    """Extract shape as element dict."""
    emu_per_px = get_emu_per_px()
    try:
        elem = {
            "type": "shape",
            "x": round(shape.left / emu_per_px),
            "y": round(shape.top / emu_per_px),
            "width": round(shape.width / emu_per_px),
            "height": round(shape.height / emu_per_px),
            "shape": _resolve_shape_name(shape)
        }
        _add_flip(elem, shape)
        
        # Style references
        style_fill_idx = None
        style_fill_color = None
        try:
            style = shape._element.find(f'{{{_NS["p"]}}}style')
            if style is not None:
                fill_ref = style.find(f'{{{_NS["a"]}}}fillRef')
                if fill_ref is not None:
                    style_fill_idx = int(fill_ref.get('idx', 0))
                    sc = fill_ref.find(f'{{{_NS["a"]}}}schemeClr')
                    if sc is not None:
                        style_fill_color = sc.get('val')
        except Exception:
            pass
        
        # Rotation
        if shape.rotation != 0:
            elem["rotation"] = round(shape.rotation, 1)
        
        # Adjustments (only if explicitly set in XML avLst)
        try:
            sp_pr = shape._element.spPr
            prst_geom = sp_pr.find(f'{{{_NS["a"]}}}prstGeom')
            if prst_geom is not None:
                av_lst = prst_geom.find(f'{{{_NS["a"]}}}avLst')
                if av_lst is not None and len(av_lst) > 0:
                    adjs = []
                    for gd in av_lst.findall(f'{{{_NS["a"]}}}gd'):
                        fmla = gd.get('fmla', '')
                        if fmla.startswith('val '):
                            adjs.append(round(int(fmla.split()[1]) / 100000, 5))
                    prst_name = prst_geom.get('prst')
                    if prst_name == 'arc' and len(adjs) >= 2:
                        # Raw adj are angles in 60000ths of a degree, but the
                        # builder's arc API is [startDeg, sweepDeg] — feeding
                        # raw values drew a 353° ring as ~40%.
                        start_deg = round(adjs[0] * 100000 / 60000, 3)
                        end_deg = round(adjs[1] * 100000 / 60000, 3)
                        sweep = round((end_deg - start_deg) % 360, 3)
                        adjs = [start_deg, sweep]
                    if adjs:
                        elem["adjustments"] = adjs
        except Exception:
            pass
        
        # Extract fill and line from XML
        sp_pr_xml = shape._element.find('.//{http://schemas.openxmlformats.org/presentationml/2006/main}spPr')
        elem.update(_extract_visual_effects(sp_pr_xml, theme_colors, color_mapping))
        
        # Fill (XML first, python-pptx API fallback for style references)
        try:
            fill_info = _extract_fill_from_xml(sp_pr_xml, theme_colors, color_mapping)
            # Check if spPr has explicit <a:noFill>
            has_explicit_no_fill = sp_pr_xml is not None and sp_pr_xml.find(f'{{{_NS["a"]}}}noFill') is not None
            if fill_info.get("fill") != "none" or "gradient" in fill_info or "patternFill" in fill_info:
                elem.update(fill_info)
            elif has_explicit_no_fill:
                elem["fill"] = "none"
            else:
                if shape.fill.type == 1:  # SOLID
                    if shape.fill.fore_color.type == 1:  # RGB
                        rgb = shape.fill.fore_color.rgb
                        elem["fill"] = f"#{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"
                    elif shape.fill.fore_color.type == 2:  # SCHEME
                        theme_color = shape.fill.fore_color.theme_color
                        if theme_colors and theme_color in theme_colors:
                            elem["fill"] = theme_colors[theme_color]
                    alpha_el = sp_pr_xml.find('.//{http://schemas.openxmlformats.org/drawingml/2006/main}alpha') if sp_pr_xml is not None else None
                    if alpha_el is not None:
                        elem["opacity"] = round(int(alpha_el.get('val', 100000)) / 1000, 1)
                elif shape.fill.type == 3:  # GRADIENT
                    try:
                        stops = []
                        for stop in shape.fill.gradient_stops:
                            s = {"position": round(stop.position, 3)}
                            if stop.color.type == 1:
                                rgb = stop.color.rgb
                                s["color"] = f"#{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"
                            elif stop.color.type == 2 and theme_colors and stop.color.theme_color in theme_colors:
                                s["color"] = theme_colors[stop.color.theme_color]
                            if "color" in s:
                                stops.append(s)
                        if stops:
                            angle = 0
                            try:
                                # python-pptx returns CCW angle; convert to CW (PowerPoint UI)
                                ccw = round(shape.fill.gradient_angle, 1)
                                angle = round((360 - ccw) % 360, 1)
                            except Exception:
                                pass
                            elem["gradient"] = {"stops": stops, "angle": angle}
                    except Exception:
                        pass
                elif shape.fill.type is None or shape.fill.type == 0 or shape.fill.type == 5:
                    # Resolve from style fillRef (unless useBgFill=1)
                    use_bg = shape._element.get('useBgFill') == '1'
                    if not use_bg and style_fill_idx and style_fill_idx > 0 and style_fill_color and theme_styles and theme_styles.get("fill"):
                        fill_idx = style_fill_idx - 1
                        if 0 <= fill_idx < len(theme_styles["fill"]):
                            from lxml import etree as _et
                            fill_xml = _et.fromstring(theme_styles["fill"][fill_idx])
                            scheme = fill_xml.find(f'.//{{{_NS["a"]}}}schemeClr')
                            if scheme is not None and scheme.get('val') == 'phClr':
                                resolved = _resolve_color_with_transforms(scheme, theme_colors, color_mapping, override_scheme=style_fill_color)
                                if resolved:
                                    elem["fill"] = resolved
                    if "fill" not in elem:
                        elem["fill"] = "none"
        except Exception:
            pass
        
        # Line (XML first, style reference fallback)
        try:
            line_info = _extract_line_from_xml(sp_pr_xml, theme_colors, color_mapping)
            ln_xml = sp_pr_xml.find('a:ln', _NS) if sp_pr_xml is not None else None
            if line_info.get("line") not in (None, "none") or "lineGradient" in line_info:
                # If lineWidth missing, try style reference
                if "lineWidth" not in line_info:
                    style_info = _resolve_line_from_style(shape, theme_colors, color_mapping, theme_styles)
                    if style_info.get("lineWidth"):
                        line_info["lineWidth"] = style_info["lineWidth"]
                elem.update(line_info)
            elif ln_xml is not None and len(ln_xml) > 0:
                elem.update(line_info)  # ln exists with noFill or explicit content
            else:
                elem.update(_resolve_line_from_style(shape, theme_colors, color_mapping, theme_styles))
            dash = extract_line_dash(shape)
            if dash:
                elem["dashStyle"] = dash
            # Arrow heads
            if ln_xml is not None:
                for attr, tag in [("headEnd", "headEnd"), ("tailEnd", "tailEnd")]:
                    el = ln_xml.find(f'{{{_NS["a"]}}}{tag}')
                    if el is not None and el.get('type') and el.get('type') != 'none':
                        elem[attr] = el.get('type')
        except Exception:
            if "line" not in elem and "lineGradient" not in elem:
                elem["line"] = "none"
        if "line" not in elem and "lineGradient" not in elem:
            elem["line"] = "none"
        
        # Extract text with styles
        if shape.has_text_frame and shape.text.strip():
            _extract_shape_text(shape, elem, theme_colors, color_mapping, builder_text_color=builder_text_color)
        
        # Extract hyperlink
        try:
            if hasattr(shape, 'click_action') and shape.click_action.hyperlink and shape.click_action.hyperlink.address:
                elem["link"] = shape.click_action.hyperlink.address
            else:
                # Remove null link
                if "link" in elem and elem["link"] is None:
                    del elem["link"]
        except Exception:
            pass
        
        # Extract visual effects
        elem.update(_extract_visual_effects(sp_pr_xml, theme_colors, color_mapping))

        # No effects in source → say so explicitly. The builder's add_shape
        # carries python-pptx's default <p:style> whose effectRef pulls the
        # theme shadow; an empty effectLst is needed to suppress it.
        if not any(k in elem for k in ("shadow", "glow", "softEdge", "reflection")):
            try:
                style_el = shape._element.find(f'{{{_NS["p"]}}}style')
                eff_ref = style_el.find(f'{{{_NS["a"]}}}effectRef') if style_el is not None else None
                if eff_ref is None or int(eff_ref.get('idx', '0') or 0) == 0:
                    elem["_noEffects"] = True
            except Exception:
                pass
        
        # Preserve lstStyle for roundtrip fidelity (non-placeholder shapes)
        _lst = _serialize_lstStyle(shape) if shape.has_text_frame else None
        if _lst:
            elem["_lstStyle"] = _lst
        
        return elem
    except Exception as e:
        print(f"Warning: Failed to extract shape details: {e}", file=sys.stderr)
        return None

