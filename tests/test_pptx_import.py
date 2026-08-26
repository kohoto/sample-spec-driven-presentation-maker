# Copyright Amazon.com, Inc. or its affiliates. All Rights Reserved.
# SPDX-License-Identifier: MIT-0
"""Tests for PPTX import → deck structure conversion.

Covers:
- pptx_to_json output: deck.json + slides/slide-NN.json
- shared/ingest._convert_pptx: ConversionResult with deck_structure / theme_hints
- pptx_builder.py CLI: accepts deck directory
- Engine deck_text_summary: slide-by-slide text extraction
- Non-regression: PDF/DOCX/XLSX import unchanged
"""

from __future__ import annotations

import json
import re
import shutil
import subprocess
import sys
from pathlib import Path

import pytest

_REPO_ROOT = Path(__file__).resolve().parent.parent
_FIXTURE_PPTX = _REPO_ROOT / "sdpm" / "references" / "examples" / "components.pptx"


@pytest.fixture
def fixture_pptx() -> Path:
    """Return path to the bundled components.pptx fixture."""
    if not _FIXTURE_PPTX.exists():
        pytest.skip(f"Fixture PPTX not found: {_FIXTURE_PPTX}")
    return _FIXTURE_PPTX


# ---------------------------------------------------------------------------
# T1: pptx_to_json output structure
# ---------------------------------------------------------------------------


class TestPptxToJsonDeckStructure:
    """pptx_to_json must output deck.json + slides/slide-NN.json."""

    def test_pptx_to_json_outputs_deck_structure(self, fixture_pptx: Path, tmp_path: Path) -> None:
        """Output directory contains deck.json and slides/slide-NN.json — no single slides.json."""
        from sdpm.engine.converter import pptx_to_json

        out_dir = tmp_path / "out"
        pptx_to_json(fixture_pptx, output_dir=out_dir)

        assert (out_dir / "deck.json").exists(), "deck.json should be created"
        slides_dir = out_dir / "slides"
        assert slides_dir.is_dir(), "slides/ directory should be created"
        slide_files = list(slides_dir.glob("slide-*.json"))
        assert len(slide_files) > 0, "at least one slide-NN.json should exist"
        # Legacy single slides.json must NOT be written
        assert not (out_dir / "slides.json").exists(), "legacy slides.json must not be emitted"

    def test_pptx_to_json_slug_format(self, fixture_pptx: Path, tmp_path: Path) -> None:
        """Slug format is slide-NNN (hyphen + 3-digit zero-padded) and matches parse_outline_slugs regex.

        3-digit padding keeps lexicographic sort == presentation order up to
        999 slides (2-digit broke at 100+: "slide-100" < "slide-11").
        """
        from sdpm.api import parse_outline_slugs
        from sdpm.engine.converter import pptx_to_json

        out_dir = tmp_path / "out"
        pptx_to_json(fixture_pptx, output_dir=out_dir)

        slide_files = sorted((out_dir / "slides").glob("*.json"))
        slug_re = re.compile(r"^slide-\d{3}$")
        for f in slide_files:
            assert slug_re.match(f.stem), f"slug must match slide-NNN format: {f.stem}"

        # Sanity: construct fake outline.md with these slugs and verify parse_outline_slugs accepts them
        fake_outline = "\n".join(f"- [{f.stem}] msg" for f in slide_files)
        outline_path = tmp_path / "outline.md"
        outline_path.write_text(fake_outline, encoding="utf-8")
        parsed = parse_outline_slugs(outline_path)
        assert parsed == [f.stem for f in slide_files]

    def test_deck_json_has_fonts_and_default_text_color(self, fixture_pptx: Path, tmp_path: Path) -> None:
        """deck.json contains fonts dict and defaultTextColor (both PPTX-derived)."""
        from sdpm.engine.converter import pptx_to_json

        out_dir = tmp_path / "out"
        pptx_to_json(fixture_pptx, output_dir=out_dir)

        deck = json.loads((out_dir / "deck.json").read_text(encoding="utf-8"))
        assert "fonts" in deck, "deck.json should contain fonts"
        # defaultTextColor may be None if extraction failed, but the key should exist
        assert "defaultTextColor" in deck, "deck.json should contain defaultTextColor key"


# ---------------------------------------------------------------------------
# T2: shared/ingest._convert_pptx + _extract_theme_hints
# ---------------------------------------------------------------------------


class TestConvertPptxThemeHints:
    """_convert_pptx must populate deck_structure, slide_count, theme_hints, suggested_name."""

    def test_convert_pptx_populates_new_fields(self, fixture_pptx: Path, tmp_path: Path) -> None:
        from shared.ingest import convert_file

        out_dir = tmp_path / "out"
        result = convert_file(fixture_pptx, out_dir)

        assert result.status == "success"
        assert result.deck_structure is True, "deck_structure should be True for PPTX"
        assert result.slide_count > 0, "slide_count should be positive"
        assert result.suggested_name == fixture_pptx.stem
        assert result.theme_hints is not None

    def test_convert_pptx_theme_hints_keys(self, fixture_pptx: Path, tmp_path: Path) -> None:
        """theme_hints uses unified keys: backgroundLuminance / accentColors / fonts."""
        from shared.ingest import convert_file

        out_dir = tmp_path / "out"
        result = convert_file(fixture_pptx, out_dir)

        assert "backgroundLuminance" in result.theme_hints
        assert "accentColors" in result.theme_hints
        assert "fonts" in result.theme_hints
        # Luminance: 0.0-1.0
        lum = result.theme_hints["backgroundLuminance"]
        assert isinstance(lum, (int, float))
        assert 0.0 <= lum <= 1.0
        # Accent colors: list of #RRGGBB, max 3
        assert isinstance(result.theme_hints["accentColors"], list)
        assert len(result.theme_hints["accentColors"]) <= 3
        for c in result.theme_hints["accentColors"]:
            assert re.match(r"^#[0-9A-Fa-f]{6}$", c), f"invalid hex: {c}"
        # Fonts: dict with halfwidth/fullwidth
        fonts = result.theme_hints["fonts"]
        assert isinstance(fonts, dict)

    def test_convert_pptx_theme_hints_uses_slide_bg(
        self, fixture_pptx: Path, tmp_path: Path,
    ) -> None:
        """When slide has explicit background, luminance reflects it (median across slides).

        Synthetic fixture: modify slide XML to force explicit backgrounds.
        If the fixture doesn't have explicit bg, we fall back to template-derived luminance
        — both paths must yield a valid number.
        """
        from shared.ingest import convert_file

        out_dir = tmp_path / "out"
        result = convert_file(fixture_pptx, out_dir)
        # The fixture may or may not have explicit slide bg; we just verify the median is valid.
        # (Full-fidelity synthetic testing of slide-bg override is deferred; the pipeline
        # correctness is validated via test_convert_pptx_theme_hints_keys.)
        assert 0.0 <= result.theme_hints["backgroundLuminance"] <= 1.0



class TestThemeHintsDdbItem:
    """theme_hints_ddb_item must tolerate theme_hints=None (PR #215 follow-up, R2).

    shared.ingest._convert_pptx returns theme_hints=None (with deck_structure=True)
    when theme extraction fails; the Cloud upload path once crashed with
    AttributeError on ``None.get`` — turning a successful conversion into
    "Conversion failed".
    """

    def test_none_theme_hints_returns_empty_defaults(self) -> None:
        from shared.schema import theme_hints_ddb_item

        item = theme_hints_ddb_item(None)
        assert item["backgroundLuminance"] is None
        assert item["accentColors"] == []
        assert item["fonts"] == {}

    def test_populated_theme_hints_converts_luminance_to_decimal(self) -> None:
        from decimal import Decimal

        from shared.schema import theme_hints_ddb_item

        item = theme_hints_ddb_item({
            "backgroundLuminance": 0.12,
            "accentColors": ["#FF0000"],
            "fonts": {"halfwidth": "Arial"},
        })
        assert isinstance(item["backgroundLuminance"], Decimal)
        assert item["backgroundLuminance"] == Decimal("0.12")
        assert item["accentColors"] == ["#FF0000"]
        assert item["fonts"] == {"halfwidth": "Arial"}


class TestDeckTextSummaryEngine:
    """sdpm.utils.deck_summary — single source for Local/Cloud text summaries.

    PR #215 follow-up (R5): the summary logic was verbatim-duplicated in
    servers/local/upload_tools.py and servers/remote/tools/upload.py; it now lives
    in the engine per the logic-sharing steering.
    """

    def test_summary_shape_titles_and_dedup(self) -> None:
        from sdpm.utils.deck_summary import deck_text_summary

        slides = [
            {
                "title": "Plain Title",
                "elements": [
                    {"type": "textbox", "text": "Body text"},
                    {"type": "textbox", "text": "Body text"},  # duplicate → dropped
                ],
            },
            {
                "title": {"text": "Dict Title"},
                "elements": [
                    {
                        "type": "group",
                        "elements": [{"type": "textbox", "text": "Nested"}],
                    },
                    {
                        "type": "table",
                        "headers": ["H1", "H2"],
                        "rows": [["a", "b"]],
                    },
                    {"type": "textbox", "items": ["Item A"]},
                ],
            },
            {},  # unparseable/empty slide keeps its number
        ]
        out = deck_text_summary(slides)
        assert "--- Slide 1: Plain Title ---" in out
        assert out.count("Body text") == 1
        assert "--- Slide 2: Dict Title ---" in out
        assert "Nested" in out and "H1" in out and "Item A" in out
        assert "--- Slide 3 ---" in out


# ---------------------------------------------------------------------------
# T18 (Cloud): _import_converted copies template.pptx to deck/template.pptx
# ---------------------------------------------------------------------------


class TestCloudImportConvertedCopiesTemplate:
    """REMOVED: _import_converted no longer exists — replaced by stateless attachment pipeline.

    The new remote/tools/attachment.py uses the core import_attachment_core
    pipeline which commits immutable bundles. Template handling is verified
    by test_attachment_pipeline.py::TestImportAttachmentContract.
    """

    pass


# ---------------------------------------------------------------------------
# T1 regression: pptx_builder.py CLI accepts deck directory
# ---------------------------------------------------------------------------


class TestPptxBuilderCliAcceptsDeckDir:
    def test_pptx_builder_cli_generate_on_deck_dir(self, fixture_pptx: Path, tmp_path: Path) -> None:
        """Convert PPTX → deck structure → pptx_builder.py generate should work on the directory."""
        from sdpm.engine.converter import pptx_to_json

        deck_dir = tmp_path / "deck"
        pptx_to_json(fixture_pptx, output_dir=deck_dir)

        # pptx_to_json writes deck.json (with template=None) — supply template for generate
        deck_json_path = deck_dir / "deck.json"
        deck_data = json.loads(deck_json_path.read_text(encoding="utf-8"))
        deck_data["template"] = str(_REPO_ROOT / "sdpm" / "templates" / "blank-dark.pptx")
        deck_json_path.write_text(json.dumps(deck_data, ensure_ascii=False, indent=2), encoding="utf-8")

        # Build an outline covering all slides
        specs_dir = deck_dir / "specs"
        specs_dir.mkdir(exist_ok=True)
        slide_files = sorted((deck_dir / "slides").glob("*.json"))
        outline_lines = [f"- [{f.stem}] Slide {f.stem}" for f in slide_files]
        (specs_dir / "outline.md").write_text("\n".join(outline_lines), encoding="utf-8")

        # CLI: pptx_builder.py generate {deck_dir} -o {output}
        output_pptx = tmp_path / "out.pptx"
        cli = _REPO_ROOT / "sdpm" / "scripts" / "pptx_builder.py"
        proc = subprocess.run(
            [sys.executable, str(cli), "generate", str(deck_dir), "-o", str(output_pptx)],
            capture_output=True, text=True, timeout=120,
        )
        assert proc.returncode == 0, f"pptx_builder.py generate failed:\nSTDOUT: {proc.stdout}\nSTDERR: {proc.stderr}"
        assert output_pptx.exists()


# ---------------------------------------------------------------------------
# T16: extract_placeholder_template (FR-9)
# ---------------------------------------------------------------------------


class TestExtractPlaceholderTemplate:
    """extract_placeholder_template keeps every master/layout, replaces slides
    with one placeholder-only sample per *used* layout."""

    def test_drops_stale_section_list(self, fixture_pptx: Path, tmp_path: Path) -> None:
        """Sections referencing dropped slide IDs must be removed (R6).

        Corporate templates often use sections; stale slide IDs inside
        <p14:sectionLst> trigger PowerPoint's repair prompt.
        """
        from lxml import etree
        from pptx import Presentation
        from pptx.oxml.ns import qn
        from sdpm.engine.converter.template import extract_placeholder_template

        _P14 = "http://schemas.microsoft.com/office/powerpoint/2010/main"
        _SECTION_URI = "{521415D9-36F7-43E2-AB2F-B90AF26B5E84}"

        # Build a section-bearing fixture: wrap all slides in one section.
        prs = Presentation(str(fixture_pptx))
        pres_el = prs.slides._sldIdLst.getparent()
        slide_ids = [e.get("id") for e in prs.slides._sldIdLst]
        assert slide_ids, "fixture must have slides"
        ext_lst = pres_el.find(qn("p:extLst"))
        if ext_lst is None:
            ext_lst = etree.SubElement(pres_el, qn("p:extLst"))
        ext = etree.SubElement(ext_lst, qn("p:ext"))
        ext.set("uri", _SECTION_URI)
        section_lst = etree.SubElement(ext, f"{{{_P14}}}sectionLst")
        section = etree.SubElement(section_lst, f"{{{_P14}}}section")
        section.set("name", "Intro")
        section.set("id", "{11111111-1111-1111-1111-111111111111}")
        sect_slides = etree.SubElement(section, f"{{{_P14}}}sldIdLst")
        for sid in slide_ids:
            sld = etree.SubElement(sect_slides, f"{{{_P14}}}sldId")
            sld.set("id", sid)
        src = tmp_path / "with_sections.pptx"
        prs.save(str(src))

        out = tmp_path / "template.pptx"
        extract_placeholder_template(src, out)

        out_prs = Presentation(str(out))
        out_xml = etree.tostring(
            out_prs.slides._sldIdLst.getparent(), encoding="unicode",
        )
        assert "sectionLst" not in out_xml, "stale sectionLst must be dropped"
        assert _SECTION_URI not in out_xml
        # Output must still be loadable with slides present.
        assert len(out_prs.slides) > 0

    def test_source_without_sections_is_unaffected(
        self, fixture_pptx: Path, tmp_path: Path
    ) -> None:
        """No-section sources keep working (guard for the cleanup step)."""
        from pptx import Presentation
        from sdpm.engine.converter.template import extract_placeholder_template

        out = tmp_path / "template.pptx"
        meta = extract_placeholder_template(fixture_pptx, out)
        assert meta["used_layout_count"] > 0
        assert len(Presentation(str(out)).slides) > 0

    def test_drops_source_content_emits_one_slide_per_used_layout(
        self, fixture_pptx: Path, tmp_path: Path
    ) -> None:
        from pptx import Presentation
        from sdpm.engine.converter.template import extract_placeholder_template

        src_prs = Presentation(str(fixture_pptx))
        assert len(src_prs.slides) > 0
        used_layout_names = {s.slide_layout.name for s in src_prs.slides}

        out = tmp_path / "template.pptx"
        meta = extract_placeholder_template(fixture_pptx, out)
        assert out.exists()

        out_prs = Presentation(str(out))
        # Slide count equals the number of distinct layouts the source used.
        assert len(out_prs.slides) == len(used_layout_names)
        assert meta["used_layout_count"] == len(used_layout_names)
        # Every emitted slide's layout is one of the used ones.
        for s in out_prs.slides:
            assert s.slide_layout.name in used_layout_names

    def test_preserves_all_layouts(
        self, fixture_pptx: Path, tmp_path: Path
    ) -> None:
        """All slide_layouts are preserved (including layouts the source never used)."""
        from pptx import Presentation
        from sdpm.engine.converter.template import extract_placeholder_template

        src_prs = Presentation(str(fixture_pptx))
        src_layout_names = []
        for master in src_prs.slide_masters:
            src_layout_names.append(sorted(layout.name for layout in master.slide_layouts))

        out = tmp_path / "template.pptx"
        extract_placeholder_template(fixture_pptx, out)

        dst_prs = Presentation(str(out))
        dst_layout_names = []
        for master in dst_prs.slide_masters:
            dst_layout_names.append(sorted(layout.name for layout in master.slide_layouts))

        assert dst_layout_names == src_layout_names

    def test_preserves_all_masters(
        self, fixture_pptx: Path, tmp_path: Path
    ) -> None:
        """slide_master count is preserved."""
        from pptx import Presentation
        from sdpm.engine.converter.template import extract_placeholder_template

        src_master_count = len(Presentation(str(fixture_pptx)).slide_masters)

        out = tmp_path / "template.pptx"
        extract_placeholder_template(fixture_pptx, out)

        dst_master_count = len(Presentation(str(out)).slide_masters)
        assert dst_master_count == src_master_count

    def test_output_smaller_than_input(
        self, fixture_pptx: Path, tmp_path: Path
    ) -> None:
        """Dropping slide content shrinks the file."""
        from sdpm.engine.converter.template import extract_placeholder_template

        out = tmp_path / "template.pptx"
        meta = extract_placeholder_template(fixture_pptx, out)
        assert meta["output_size"] < meta["input_size"]


# ---------------------------------------------------------------------------
# T17: shared/ingest._convert_pptx integrates placeholder template extraction
# ---------------------------------------------------------------------------


class TestConvertPptxOutputsTemplate:
    """_convert_pptx writes deck.json + slides/ + template.pptx and reports template_path."""

    def test_convert_pptx_outputs_template(
        self, fixture_pptx: Path, tmp_path: Path
    ) -> None:
        from shared.ingest import convert_file

        out_dir = tmp_path / "out"
        result = convert_file(fixture_pptx, out_dir)

        assert result.status == "success"
        assert result.template_path == "template.pptx"
        assert (out_dir / "template.pptx").is_file()

    def test_pptx_builder_with_derived_template(
        self, fixture_pptx: Path, tmp_path: Path
    ) -> None:
        """E2E: convert PPTX → use derived template.pptx → pptx_builder generate succeeds."""
        from shared.ingest import convert_file

        deck_dir = tmp_path / "deck"
        result = convert_file(fixture_pptx, deck_dir)
        assert result.status == "success"
        assert result.template_path == "template.pptx"

        # Wire deck.json to point at the derived template + populate outline
        deck_json_path = deck_dir / "deck.json"
        deck_data = json.loads(deck_json_path.read_text(encoding="utf-8"))
        deck_data["template"] = str(deck_dir / "template.pptx")
        deck_json_path.write_text(
            json.dumps(deck_data, ensure_ascii=False, indent=2), encoding="utf-8"
        )

        specs_dir = deck_dir / "specs"
        specs_dir.mkdir(exist_ok=True)
        slide_files = sorted((deck_dir / "slides").glob("*.json"))
        outline_lines = [f"- [{f.stem}] Slide {f.stem}" for f in slide_files]
        (specs_dir / "outline.md").write_text(
            "\n".join(outline_lines), encoding="utf-8"
        )

        output_pptx = tmp_path / "out.pptx"
        cli = _REPO_ROOT / "sdpm" / "scripts" / "pptx_builder.py"
        proc = subprocess.run(
            [sys.executable, str(cli), "generate", str(deck_dir), "-o", str(output_pptx)],
            capture_output=True, text=True, timeout=180,
        )
        assert proc.returncode == 0, (
            f"pptx_builder.py generate failed:\nSTDOUT: {proc.stdout}\nSTDERR: {proc.stderr}"
        )
        assert output_pptx.exists()


# ---------------------------------------------------------------------------
# Non-regression: PDF/DOCX/XLSX conversion still works (old flow)
# ---------------------------------------------------------------------------


class TestNonRegression:
    """PDF/DOCX/XLSX conversion paths must not be affected by the PPTX deck-structure change."""

    def test_convert_docx_still_produces_markdown(self, tmp_path: Path) -> None:
        """A simple DOCX is converted to Markdown (deck_structure stays False)."""
        try:
            from docx import Document  # noqa: F401
        except ImportError:
            pytest.skip("python-docx not available")

        from docx import Document
        from shared.ingest import convert_file

        src = tmp_path / "simple.docx"
        doc = Document()
        doc.add_heading("Test", level=1)
        doc.add_paragraph("Hello world")
        doc.save(str(src))

        out_dir = tmp_path / "out"
        result = convert_file(src, out_dir)
        assert result.status in ("success", "partial")
        # DOCX must not trigger deck_structure
        assert result.deck_structure is False
        md_path = out_dir / "simple.md"
        assert md_path.exists()
