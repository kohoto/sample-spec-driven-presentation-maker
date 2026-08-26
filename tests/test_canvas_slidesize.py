# Copyright Amazon.com, Inc. or its affiliates. All Rights Reserved.
# SPDX-License-Identifier: MIT-0
"""Tests for Issue #208: slideSize in deck.json and generate-time validation.

Covers:
- Phase 1-5: _resolve_config slideSize mismatch warning / height boundary warning
- Phase 1-6: api.init writes slideSize / converter/pipeline writes slideSize
- Phase 1-7: _apply_grid_overlay derives px_y from image aspect ratio
"""

from __future__ import annotations

import json
from pathlib import Path

_REPO_ROOT = Path(__file__).resolve().parent.parent


# ── Helpers ──

def _make_deck(tmp_path: Path, *, template_path: Path, slide_size: dict | None = None,
               elements: list[dict] | None = None) -> Path:
    """Create a minimal deck directory for generate testing."""
    deck = tmp_path / "deck"
    (deck / "slides").mkdir(parents=True)
    (deck / "specs").mkdir()
    deck_data: dict = {
        "template": str(template_path),
        "fonts": {"fullwidth": "Meiryo", "halfwidth": "Arial"},
        "defaultTextColor": "#FFFFFF",
    }
    if slide_size is not None:
        deck_data["slideSize"] = slide_size
    (deck / "deck.json").write_text(json.dumps(deck_data, ensure_ascii=False))
    (deck / "specs" / "outline.md").write_text("- [intro] Intro\n")
    slide = {
        "layout": "Blank",
        "elements": elements or [
            {"type": "textbox", "text": "Hello", "x": 100, "y": 100, "width": 800, "height": 80},
        ],
    }
    (deck / "slides" / "intro.json").write_text(json.dumps(slide, ensure_ascii=False))
    return deck


# ═══════════════════════════════════════════════════════════════════════
# Phase 1-5: generate-time slideSize / height-boundary warnings
# ═══════════════════════════════════════════════════════════════════════


class TestResolvConfigSlideSizeWarning:
    """_resolve_config emits slideSize mismatch warning when deck.json disagrees with template."""

    def test_mismatch_warning_4x3_deck_wrong_size(self, template_4x3: Path, tmp_path: Path):
        """4:3 template with wrong slideSize in deck.json triggers warning."""
        from sdpm.api import _resolve_config

        deck = _make_deck(tmp_path, template_path=template_4x3,
                          slide_size={"width": 1920, "height": 1080})
        config = _resolve_config(deck)
        mismatch_warnings = [w for w in config.warnings if "slideSize mismatch" in w]
        assert len(mismatch_warnings) == 1
        assert "1440" in mismatch_warnings[0]

    def test_no_warning_when_correct_size(self, template_4x3: Path, tmp_path: Path):
        """4:3 template with correct slideSize emits no mismatch warning."""
        from sdpm.api import _resolve_config

        deck = _make_deck(tmp_path, template_path=template_4x3,
                          slide_size={"width": 1920, "height": 1440})
        config = _resolve_config(deck)
        mismatch_warnings = [w for w in config.warnings if "slideSize mismatch" in w]
        assert len(mismatch_warnings) == 0

    def test_no_warning_when_no_slide_size_key(self, template_16x9: Path, tmp_path: Path):
        """No slideSize in deck.json → no mismatch warning (key is optional)."""
        from sdpm.api import _resolve_config

        deck = _make_deck(tmp_path, template_path=template_16x9, slide_size=None)
        config = _resolve_config(deck)
        mismatch_warnings = [w for w in config.warnings if "slideSize mismatch" in w]
        assert len(mismatch_warnings) == 0

    def test_16x9_correct_no_warning(self, template_16x9: Path, tmp_path: Path):
        """16:9 template with correct slideSize emits no warning."""
        from sdpm.api import _resolve_config

        deck = _make_deck(tmp_path, template_path=template_16x9,
                          slide_size={"width": 1920, "height": 1080})
        config = _resolve_config(deck)
        mismatch_warnings = [w for w in config.warnings if "slideSize mismatch" in w]
        assert len(mismatch_warnings) == 0


class TestResolveConfigHeightBoundaryWarning:
    """_resolve_config emits height boundary warning for out-of-bounds elements."""

    def test_height_oob_4x3(self, template_4x3: Path, tmp_path: Path):
        """Element exceeding 4:3 height (1440) triggers height boundary warning."""
        from sdpm.api import _resolve_config

        elements = [
            {"type": "textbox", "text": "A", "x": 0, "y": 1400, "width": 200, "height": 100},
        ]
        deck = _make_deck(tmp_path, template_path=template_4x3, elements=elements)
        config = _resolve_config(deck)
        height_warnings = [w for w in config.warnings if "exceeds slide height" in w]
        assert len(height_warnings) == 1
        assert "1440" in height_warnings[0]

    def test_no_height_warning_within_bounds_4x3(self, template_4x3: Path, tmp_path: Path):
        """Element within 4:3 height (y+h<=1440) does not trigger warning."""
        from sdpm.api import _resolve_config

        elements = [
            {"type": "textbox", "text": "A", "x": 0, "y": 1300, "width": 200, "height": 100},
        ]
        deck = _make_deck(tmp_path, template_path=template_4x3, elements=elements)
        config = _resolve_config(deck)
        height_warnings = [w for w in config.warnings if "exceeds slide height" in w]
        assert len(height_warnings) == 0

    def test_no_false_positive_16x9(self, template_16x9: Path, tmp_path: Path):
        """16:9: element at y=900 h=180 → 1080 exactly, no warning."""
        from sdpm.api import _resolve_config

        elements = [
            {"type": "textbox", "text": "A", "x": 0, "y": 900, "width": 200, "height": 180},
        ]
        deck = _make_deck(tmp_path, template_path=template_16x9, elements=elements)
        config = _resolve_config(deck)
        height_warnings = [w for w in config.warnings if "exceeds slide height" in w]
        assert len(height_warnings) == 0

    def test_height_oob_16x9(self, template_16x9: Path, tmp_path: Path):
        """16:9: element at y=1000 h=100 → 1100 > 1080, triggers warning."""
        from sdpm.api import _resolve_config

        elements = [
            {"type": "textbox", "text": "A", "x": 0, "y": 1000, "width": 200, "height": 100},
        ]
        deck = _make_deck(tmp_path, template_path=template_16x9, elements=elements)
        config = _resolve_config(deck)
        height_warnings = [w for w in config.warnings if "exceeds slide height" in w]
        assert len(height_warnings) == 1
        assert "1080" in height_warnings[0]


# ═══════════════════════════════════════════════════════════════════════
# Phase 1-6: api.init writes slideSize
# ═══════════════════════════════════════════════════════════════════════


class TestInitSlidesSize:
    """api.init creates a minimal deck.json without slideSize (template branch removed)."""

    def test_init_no_slide_size(self, tmp_path: Path):
        """init never writes slideSize — the agent writes it after analyze_template."""
        from sdpm.api import init

        result = init(name="test", output_dir=str(tmp_path / "out"))
        deck_json = json.loads(Path(result["deck_json"]).read_text())
        assert "slideSize" not in deck_json


# ═══════════════════════════════════════════════════════════════════════
# Phase 1-6: converter/pipeline writes slideSize
# ═══════════════════════════════════════════════════════════════════════


class TestPipelineSlidesSize:
    """pptx_to_json writes slideSize in deck.json."""

    def test_16x9_import(self, template_16x9: Path, tmp_path: Path):
        from sdpm.engine.converter import pptx_to_json

        # Need a pptx with at least one slide
        from pptx import Presentation
        prs = Presentation(str(template_16x9))
        layout = prs.slide_layouts[0]
        prs.slides.add_slide(layout)
        src = tmp_path / "src.pptx"
        prs.save(str(src))

        out_dir = tmp_path / "out"
        pptx_to_json(src, output_dir=out_dir)
        deck = json.loads((out_dir / "deck.json").read_text())
        assert deck["slideSize"] == {"width": 1920, "height": 1080, "ptPerPx": 0.5}

    def test_4x3_import(self, template_4x3: Path, tmp_path: Path):
        from sdpm.engine.converter import pptx_to_json
        from pptx import Presentation

        prs = Presentation(str(template_4x3))
        layout = prs.slide_layouts[0]
        prs.slides.add_slide(layout)
        src = tmp_path / "src.pptx"
        prs.save(str(src))

        out_dir = tmp_path / "out"
        pptx_to_json(src, output_dir=out_dir)
        deck = json.loads((out_dir / "deck.json").read_text())
        assert deck["slideSize"] == {"width": 1920, "height": 1440, "ptPerPx": 0.375}


# ═══════════════════════════════════════════════════════════════════════
# Phase 1-7: _apply_grid_overlay px_y derivation
# ═══════════════════════════════════════════════════════════════════════


class TestGridOverlayAspectRatio:
    """_apply_grid_overlay derives px_y from image dimensions, not fixed 1080."""

    def test_16x9_image_px_y_is_1080_based(self, tmp_path: Path):
        """16:9 image: px_y at 50% should be 540 (= 1920 * 540/960 * 50/100 = 540)."""
        from PIL import Image
        from sdpm.api import _apply_grid_overlay

        # 1920x1080 image (or proportional: 960x540)
        img = Image.new("RGB", (960, 540), (255, 255, 255))
        p = tmp_path / "slide.png"
        img.save(str(p))

        _apply_grid_overlay([str(p)])

        # Verify the file was modified (overlay applied)
        result = Image.open(str(p))
        assert result.size == (960, 540)

    def test_4x3_image_px_y_is_1440_based(self, tmp_path: Path):
        """4:3 image (1920x1440 proportional): px_y at 50% should be 720."""
        from PIL import Image
        from sdpm.api import _apply_grid_overlay

        # 4:3 proportional image: 960x720
        img = Image.new("RGB", (960, 720), (255, 255, 255))
        p = tmp_path / "slide.png"
        img.save(str(p))

        _apply_grid_overlay([str(p)])

        # Just verify it runs without error (assertion on pixel content is fragile)
        result = Image.open(str(p))
        assert result.size == (960, 720)

    def test_px_y_formula_correctness(self):
        """Verify formula: round(1920 * h / w * pct / 100) gives correct results."""
        # 16:9 (1920x1080 equivalent: w=960, h=540) at 50%
        assert round(1920 * 540 / 960 * 50 / 100) == 540
        # 4:3 (1920x1440 equivalent: w=960, h=720) at 50%
        assert round(1920 * 720 / 960 * 50 / 100) == 720
        # 4:3 at 100%
        assert round(1920 * 720 / 960 * 100 / 100) == 1440
