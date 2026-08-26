# Copyright Amazon.com, Inc. or its affiliates. All Rights Reserved.
# SPDX-License-Identifier: MIT-0
"""Conversion-scale isolation tests (v0.5.2).

The converter's px scale (EMU per px on the 1920px basis) used to be a
process-global mutated through ``sys.modules`` patching — a silent-regression
trap for non-standard slide widths and concurrent conversions. It is now a
``ContextVar`` scoped by ``conversion_scale``. These tests pin:

- default value and scoped set/restore (normal, exception, nested)
- thread isolation (two concurrent conversions with different widths)
- pipeline conversions never leak scale into later direct extractor calls
- AST guards: the legacy mechanisms must not be reintroduced
"""

import ast
import threading
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Emu

from sdpm.engine.converter.constants import conversion_scale, get_emu_per_px

_DEFAULT = 6350.0
_W_16X9 = 12192000  # -> 6350.0
_W_4X3 = 9144000    # -> 4762.5


class TestScopedScale:
    def test_default_scale(self):
        assert get_emu_per_px() == _DEFAULT

    def test_scope_sets_and_restores(self):
        with conversion_scale(_W_4X3):
            assert get_emu_per_px() == 4762.5
        assert get_emu_per_px() == _DEFAULT

    def test_scope_restores_on_exception(self):
        with pytest.raises(RuntimeError):
            with conversion_scale(_W_4X3):
                raise RuntimeError("boom")
        assert get_emu_per_px() == _DEFAULT

    def test_nested_scopes_restore_outer(self):
        with conversion_scale(_W_4X3):
            assert get_emu_per_px() == 4762.5
            with conversion_scale(_W_16X9):
                assert get_emu_per_px() == 6350.0
            assert get_emu_per_px() == 4762.5
        assert get_emu_per_px() == _DEFAULT

    def test_threads_are_isolated(self):
        """Two conversions with different widths must not see each other's scale."""
        barrier = threading.Barrier(2, timeout=10)
        results = {}
        errors = []

        def worker(idx, width):
            try:
                with conversion_scale(width):
                    barrier.wait()  # both threads are inside their scopes now
                    results[idx] = get_emu_per_px()
                    barrier.wait()  # hold the scope until both have read
            except Exception as e:  # pragma: no cover - failure reporting
                errors.append(e)

        threads = [
            threading.Thread(target=worker, args=(0, _W_16X9)),
            threading.Thread(target=worker, args=(1, _W_4X3)),
        ]
        for t in threads:
            t.start()
        for t in threads:
            t.join(timeout=10)
        assert not errors
        assert results == {0: 6350.0, 1: 4762.5}


class TestPipelineScaleIsolation:
    def _make_4x3_deck(self, tmp_path):
        prs = Presentation()
        prs.slide_width = Emu(_W_4X3)
        prs.slide_height = Emu(6858000)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        emu43 = _W_4X3 / 1920
        tb = slide.shapes.add_textbox(
            Emu(round(100 * emu43)), Emu(round(50 * emu43)),
            Emu(round(300 * emu43)), Emu(round(80 * emu43)))
        tb.text_frame.text = "43"
        path = tmp_path / "deck43.pptx"
        prs.save(path)
        return path

    def test_non_standard_pipeline_does_not_leak(self, tmp_path):
        """After a 4:3 conversion, direct extractor calls are back on the 16:9 basis."""
        from sdpm.engine.converter.pipeline import pptx_to_json
        from sdpm.engine.converter.elements import extract_textbox_element

        result = pptx_to_json(self._make_4x3_deck(tmp_path), tmp_path / "out")
        elem43 = result["slides"][0]["elements"][0]
        assert (elem43["x"], elem43["y"], elem43["width"]) == (100, 50, 300)

        # Scale must be fully restored...
        assert get_emu_per_px() == _DEFAULT
        # ...so a direct extraction now uses the default 16:9 basis.
        prs = Presentation()
        prs.slide_width = Emu(_W_16X9)
        prs.slide_height = Emu(6858000)
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tb = slide.shapes.add_textbox(Emu(200 * 6350), Emu(100 * 6350), Emu(400 * 6350), Emu(50 * 6350))
        tb.text_frame.text = "after"
        elem = extract_textbox_element(tb)
        assert (elem["x"], elem["y"], elem["width"]) == (200, 100, 400)

    def test_pipeline_inside_outer_scope_restores_outer(self, tmp_path):
        """A conversion nested inside another scale scope restores the outer scale."""
        from sdpm.engine.converter.pipeline import pptx_to_json

        deck = self._make_4x3_deck(tmp_path)
        with conversion_scale(_W_16X9):
            result = pptx_to_json(deck, tmp_path / "out_nested")
            assert result["slides"][0]["elements"][0]["x"] == 100  # inner deck's own basis
            assert get_emu_per_px() == 6350.0  # outer scope intact
        assert get_emu_per_px() == _DEFAULT


# ---------------------------------------------------------------------------
# AST guards: the legacy scale mechanisms must never come back
# ---------------------------------------------------------------------------

_CONVERTER_DIR = Path(__file__).resolve().parents[1] / "sdpm" / "sdpm" / "engine" / "converter"
_CONVERTER_FILES = sorted(_CONVERTER_DIR.rglob("*.py"))


def _is_facade(path: Path) -> bool:
    """The package facade may re-export the EMU_PER_PX compat constant."""
    return path.name == "__init__.py" or path.name == "constants.py"


def test_converter_files_discovered():
    assert len(_CONVERTER_FILES) >= 10


@pytest.mark.parametrize("path", _CONVERTER_FILES, ids=lambda p: str(p.relative_to(_CONVERTER_DIR)))
def test_no_static_emu_import_in_converter_internals(path):
    """Internals must use get_emu_per_px(); an import-time EMU_PER_PX copy would
    freeze the default scale and silently break non-standard slide widths."""
    if _is_facade(path):
        pytest.skip("facade/constants keep the compat constant")
    tree = ast.parse(path.read_text())
    for node in ast.walk(tree):
        if isinstance(node, ast.ImportFrom):
            names = [a.name for a in node.names]
            assert "EMU_PER_PX" not in names, (
                f"{path.name} line {node.lineno}: import EMU_PER_PX is forbidden — "
                "use get_emu_per_px() (scoped scale) instead")


@pytest.mark.parametrize("path", _CONVERTER_FILES, ids=lambda p: str(p.relative_to(_CONVERTER_DIR)))
def test_no_sys_modules_patching_in_converter(path):
    """The sys.modules module-name patch list (old set_emu_per_px) must not return."""
    tree = ast.parse(path.read_text())
    for node in ast.walk(tree):
        if (isinstance(node, ast.Attribute) and node.attr == "modules"
                and isinstance(node.value, ast.Name) and node.value.id == "sys"):
            raise AssertionError(
                f"{path.name} line {node.lineno}: sys.modules access is forbidden in converter")
        if isinstance(node, (ast.FunctionDef, ast.AsyncFunctionDef)):
            assert node.name != "set_emu_per_px", (
                f"{path.name} line {node.lineno}: set_emu_per_px must not be reintroduced")
