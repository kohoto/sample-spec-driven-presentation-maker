# Copyright Amazon.com, Inc. or its affiliates. All Rights Reserved.
# SPDX-License-Identifier: MIT-0
"""Characterization tests for converter element extraction (v0.5.2 split safety net).

These tests pin the exact observable behavior of the element extractors and the
conversion pipeline *before* ``converter/elements.py`` is split into a package.

RULES (from the v0.5.2 SPEC):
- Expected values below were captured from the pre-split implementation.
  DO NOT update them to make a refactoring pass — a mismatch means the
  refactoring changed behavior and must be reverted.
- Fixtures are generated programmatically with python-pptx (no opaque binary
  fixtures in the repo).

Covers:
- Direct output contract of the 7 public extractors (+ ``_dispatch_shape``)
- Image/media file naming and the image counter
- Pipeline JSON for a standard 16:9 deck and a non-standard 4:3 deck
  (coordinates always normalize to the 1920px basis)
- Import-path and signature compatibility for the surface other modules use
"""

import inspect
import io
import json

import pytest
from lxml import etree
from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.util import Emu

from sdpm.engine.converter.elements import (
    _dispatch_shape,
    extract_freeform_element,
    extract_group_element,
    extract_line_element,
    extract_picture_element,
    extract_shape_element,
    extract_textbox_element,
    extract_video_element,
)

# EMU per px on the 1920px basis for a standard 16:9 deck (12192000 / 1920)
EMU = 6350

_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
_P = "http://schemas.openxmlformats.org/presentationml/2006/main"
_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
_ASVG = "http://schemas.microsoft.com/office/drawing/2016/SVG/main"
_P14 = "http://schemas.microsoft.com/office/powerpoint/2010/main"


@pytest.fixture(autouse=True)
def _restore_converter_scale():
    """Undo any conversion-scale leak between tests.

    The legacy implementation mutates module globals via ``set_emu_per_px``;
    a 4:3 pipeline run would poison later direct-extractor tests. Once the
    scoped-scale implementation lands this becomes a no-op.
    """
    yield
    import sdpm.engine.converter.constants as _c

    if hasattr(_c, "set_emu_per_px"):  # legacy global-state implementation
        _c.set_emu_per_px(1920 * 6350)


def _blank_slide(width_emu=12192000, height_emu=6858000):
    prs = Presentation()
    prs.slide_width = Emu(width_emu)
    prs.slide_height = Emu(height_emu)
    return prs, prs.slides.add_slide(prs.slide_layouts[6])


def _png_stream(size=(16, 12), color=(200, 30, 30)):
    buf = io.BytesIO()
    Image.new("RGB", size, color).save(buf, format="PNG")
    buf.seek(0)
    return buf


# ---------------------------------------------------------------------------
# Direct extractor contracts
# ---------------------------------------------------------------------------


class TestLineExtractor:
    def test_straight_connector(self):
        _, slide = _blank_slide()
        conn = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, Emu(100 * EMU), Emu(200 * EMU), Emu(500 * EMU), Emu(400 * EMU))
        assert extract_line_element(conn) == {
            "type": "line", "x1": 100, "y1": 200, "x2": 500, "y2": 400,
            "preset": "line", "color": "none",
        }

    def test_flip_and_rotation_baked_into_endpoints(self):
        """flipH swaps x endpoints; rot is baked into coordinates (schema has no line rotation)."""
        _, slide = _blank_slide()
        conn = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, Emu(100 * EMU), Emu(100 * EMU), Emu(300 * EMU), Emu(200 * EMU))
        xfrm = conn._element.spPr.find(f".//{{{_A}}}xfrm")
        xfrm.set("flipH", "1")
        xfrm.set("rot", str(90 * 60000))
        assert extract_line_element(conn) == {
            "type": "line", "x1": 250, "y1": 250, "x2": 150, "y2": 50,
            "preset": "line", "color": "none",
        }


class TestShapeExtractor:
    def test_rounded_rectangle_with_text(self):
        _, slide = _blank_slide()
        sh = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Emu(100 * EMU), Emu(100 * EMU), Emu(300 * EMU), Emu(150 * EMU))
        sh.text_frame.text = "Hello"
        assert extract_shape_element(sh) == {
            "type": "shape", "x": 100, "y": 100, "width": 300, "height": 150,
            "shape": "rounded_rectangle", "fill": "none", "line": "#41B3FF",
            "lineWidth": 0.5, "verticalAlign": "middle", "text": "Hello",
            "fontSize": 18, "align": "left",
        }

    def test_rotation(self):
        _, slide = _blank_slide()
        sh = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Emu(100 * EMU), Emu(100 * EMU), Emu(200 * EMU), Emu(100 * EMU))
        sh.rotation = 45
        assert extract_shape_element(sh) == {
            "type": "shape", "x": 100, "y": 100, "width": 200, "height": 100,
            "shape": "rectangle", "rotation": 45.0, "fill": "none",
            "line": "#41B3FF", "lineWidth": 0.5,
        }

    def test_equal_sided_oval_reports_circle(self):
        _, slide = _blank_slide()
        sh = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Emu(0), Emu(0), Emu(80 * EMU), Emu(80 * EMU))
        assert extract_shape_element(sh)["shape"] == "circle"


class TestTextboxExtractor:
    def test_single_paragraph(self):
        _, slide = _blank_slide()
        tb = slide.shapes.add_textbox(Emu(50 * EMU), Emu(60 * EMU), Emu(400 * EMU), Emu(80 * EMU))
        tb.text_frame.text = "Single line"
        assert extract_textbox_element(tb) == {
            "type": "textbox", "x": 50, "y": 60, "width": 400, "height": 80,
            "autoWidth": True, "marginLeft": 14, "marginRight": 14,
            "fill": "none", "line": "none", "text": "Single line",
        }

    def test_multi_paragraph_with_level(self):
        _, slide = _blank_slide()
        tb = slide.shapes.add_textbox(Emu(50 * EMU), Emu(200 * EMU), Emu(400 * EMU), Emu(160 * EMU))
        tf = tb.text_frame
        tf.text = "First"
        p2 = tf.add_paragraph()
        p2.text = "Second"
        p3 = tf.add_paragraph()
        p3.text = "Third"
        p3.level = 1
        assert extract_textbox_element(tb) == {
            "type": "textbox", "x": 50, "y": 200, "width": 400, "height": 160,
            "autoWidth": True, "marginLeft": 14, "marginRight": 14,
            "fill": "none", "line": "none",
            "paragraphs": [{"text": "First"}, {"text": "Second"}, {"text": "Third", "level": 1}],
        }

    def test_preset_geometry_textbox_delegates_to_shape(self):
        """A 'textbox' whose spPr carries non-rect preset geometry is really a shape."""
        _, slide = _blank_slide()
        tb = slide.shapes.add_textbox(Emu(0), Emu(0), Emu(100 * EMU), Emu(50 * EMU))
        tb.text_frame.text = "shaped"
        prst = tb._element.spPr.find(f"{{{_A}}}prstGeom")
        prst.set("prst", "ellipse")
        assert extract_textbox_element(tb)["type"] == "shape"


class TestFreeformExtractor:
    def test_triangle_path(self):
        _, slide = _blank_slide()
        fb = slide.shapes.build_freeform(Emu(0), Emu(0), scale=1.0)
        fb.add_line_segments(
            [(Emu(100 * EMU), Emu(0)), (Emu(100 * EMU), Emu(100 * EMU)), (Emu(0), Emu(100 * EMU))],
            close=True)
        ff = fb.convert_to_shape()
        expected_pathlst = (
            '<a:pathLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
            'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
            'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">'
            '<a:path w="635000" h="635000"><a:moveTo><a:pt x="0" y="0"/></a:moveTo>'
            '<a:lnTo><a:pt x="635000" y="0"/></a:lnTo>'
            '<a:lnTo><a:pt x="635000" y="635000"/></a:lnTo>'
            '<a:lnTo><a:pt x="0" y="635000"/></a:lnTo><a:close/></a:path></a:pathLst>'
        )
        assert extract_freeform_element(ff) == {
            "type": "freeform", "x": 0, "y": 0, "width": 100, "height": 100,
            "_widthEmu": 635000, "_heightEmu": 635000,
            "path": [
                {"cmd": "M", "x": 0.0, "y": 0.0},
                {"cmd": "L", "x": 100.0, "y": 0.0},
                {"cmd": "L", "x": 100.0, "y": 100.0},
                {"cmd": "L", "x": 0.0, "y": 100.0},
                {"cmd": "Z"},
            ],
            "_pathLstXml": expected_pathlst,
            "fill": "none", "line": "none",
        }


class TestPictureExtractor:
    def test_raster_picture_saved_and_named(self, tmp_path):
        _, slide = _blank_slide()
        pic = slide.shapes.add_picture(
            _png_stream(size=(64, 48)), Emu(700 * EMU), Emu(100 * EMU), Emu(128 * EMU), Emu(96 * EMU))
        elem = extract_picture_element(pic, output_dir=tmp_path, slide_idx=0, img_idx=0)
        assert elem == {
            "type": "image", "x": 700, "y": 100, "width": 128, "height": 96,
            "src": "images/slide1_image1.png", "link": None, "fit": "stretch",
        }
        assert (tmp_path / "images" / "slide1_image1.png").exists()

    def test_svg_picture(self, tmp_path):
        _, slide = _blank_slide()
        pic = slide.shapes.add_picture(
            _png_stream(size=(10, 10)), Emu(100 * EMU), Emu(100 * EMU), Emu(200 * EMU), Emu(200 * EMU))
        svg_bytes = (b'<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 100 100">'
                     b'<rect width="100" height="100" fill="green"/></svg>')
        from pptx.opc.package import Part
        from pptx.opc.packuri import PackURI
        part = Part(PackURI("/ppt/media/image_char.svg"), "image/svg+xml",
                    slide.part.package, svg_bytes)
        rid = slide.part.relate_to(part, f"{_R}/image")
        blip = pic._element.find(f".//{{{_A}}}blip")
        ext_lst = etree.SubElement(blip, f"{{{_A}}}extLst")
        ext = etree.SubElement(ext_lst, f"{{{_A}}}ext")
        ext.set("uri", "{96DAC541-7B7A-43D3-8B79-37D633B846F1}")
        svg_blip = etree.SubElement(ext, f"{{{_ASVG}}}svgBlip")
        svg_blip.set(f"{{{_R}}}embed", rid)

        elem = extract_picture_element(pic, output_dir=tmp_path, slide_idx=0, img_idx=0)
        assert elem == {
            "type": "image", "x": 100, "y": 100, "width": 200, "height": 200,
            "src": "images/slide1_image1.svg", "iconColor": "none", "fit": "stretch",
        }
        assert (tmp_path / "images" / "slide1_image1.svg").read_bytes() == svg_bytes


class _FakeRel:
    def __init__(self, target_part=None, target_ref=None):
        self.target_part = target_part
        self.target_ref = target_ref


class _FakePart:
    def __init__(self, blob, content_type="video/mp4"):
        self.blob = blob
        self.content_type = content_type


class _FakeVideoPart:
    rels = {
        "rId10": _FakeRel(target_ref="media/movie.mp4"),
        "rId11": _FakeRel(target_part=_FakePart(b"FAKE_MP4_BYTES")),
        "rId12": _FakeRel(target_part=_FakePart(b"FAKE_POSTER", "image/jpeg")),
    }


class _FakeVideoShape:
    """Minimal protocol fake — python-pptx has no video authoring API."""

    shape_type = 16  # MSO_SHAPE_TYPE.MEDIA
    left, top, width, height = 100 * EMU, 200 * EMU, 320 * EMU, 180 * EMU
    rotation = 0
    part = _FakeVideoPart()

    def __init__(self):
        self._element = etree.fromstring(
            f'<p:pic xmlns:p="{_P}" xmlns:a="{_A}" xmlns:r="{_R}" xmlns:p14="{_P14}">'
            '<p:nvPicPr><p:cNvPr id="9" name="video"/><p:cNvPicPr/>'
            '<p:nvPr><a:videoFile r:link="rId10"/>'
            '<p:extLst><p:ext uri="x"><p14:media r:embed="rId11"/></p:ext></p:extLst>'
            "</p:nvPr></p:nvPicPr>"
            '<p:blipFill><a:blip r:embed="rId12"/><a:stretch/></p:blipFill>'
            f'<p:spPr><a:xfrm><a:off x="{100 * EMU}" y="{200 * EMU}"/>'
            f'<a:ext cx="{320 * EMU}" cy="{180 * EMU}"/></a:xfrm></p:spPr></p:pic>'
        )


class TestVideoExtractor:
    def test_video_file_and_poster_saved(self, tmp_path):
        elem = extract_video_element(_FakeVideoShape(), output_dir=tmp_path, slide_idx=0, img_idx=0)
        assert elem == {
            "type": "video", "x": 100, "y": 200, "width": 320, "height": 180,
            "src": "media/slide1_video1.mp4", "poster": "images/slide1_poster1.jpg",
        }
        assert (tmp_path / "media" / "slide1_video1.mp4").read_bytes() == b"FAKE_MP4_BYTES"
        assert (tmp_path / "images" / "slide1_poster1.jpg").read_bytes() == b"FAKE_POSTER"


class TestGroupExtractor:
    def test_flat_group_children_in_slide_coordinates(self):
        _, slide = _blank_slide()
        gr = slide.shapes.add_group_shape()
        gr.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Emu(1000 * EMU), Emu(500 * EMU), Emu(100 * EMU), Emu(50 * EMU))
        gr.shapes.add_shape(
            MSO_SHAPE.OVAL, Emu(1150 * EMU), Emu(500 * EMU), Emu(80 * EMU), Emu(80 * EMU))
        elem, counter = extract_group_element(gr)
        assert counter == 0
        assert elem == {
            "type": "group", "x": 1000, "y": 500, "width": 230, "height": 80,
            "elements": [
                {"type": "shape", "x": 1000, "y": 500, "width": 100, "height": 50,
                 "shape": "rectangle", "fill": "none", "line": "#41B3FF", "lineWidth": 0.5},
                {"type": "shape", "x": 1150, "y": 500, "width": 80, "height": 80,
                 "shape": "circle", "fill": "none", "line": "#41B3FF", "lineWidth": 0.5},
            ],
        }


# ---------------------------------------------------------------------------
# Dispatch routing + image counter
# ---------------------------------------------------------------------------


class TestDispatch:
    def test_autoshape_routes_to_shape(self):
        _, slide = _blank_slide()
        sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), Emu(100 * EMU), Emu(50 * EMU))
        elem, counter = _dispatch_shape(sh)
        assert elem["type"] == "shape"
        assert counter == 0

    def test_textbox_routes_to_textbox(self):
        _, slide = _blank_slide()
        tb = slide.shapes.add_textbox(Emu(0), Emu(0), Emu(100 * EMU), Emu(50 * EMU))
        tb.text_frame.text = "t"
        elem, counter = _dispatch_shape(tb)
        assert elem["type"] == "textbox"
        assert counter == 0

    def test_line_routes_to_line(self):
        _, slide = _blank_slide()
        conn = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, Emu(0), Emu(0), Emu(100 * EMU), Emu(50 * EMU))
        elem, _ = _dispatch_shape(conn)
        assert elem["type"] == "line"

    def test_freeform_routes_to_freeform(self):
        _, slide = _blank_slide()
        fb = slide.shapes.build_freeform(Emu(0), Emu(0), scale=1.0)
        fb.add_line_segments(
            [(Emu(50 * EMU), Emu(0)), (Emu(50 * EMU), Emu(50 * EMU)), (Emu(0), Emu(50 * EMU))],
            close=True)
        ff = fb.convert_to_shape()
        elem, _ = _dispatch_shape(ff)
        assert elem["type"] == "freeform"

    def test_group_routes_recursively(self):
        _, slide = _blank_slide()
        gr = slide.shapes.add_group_shape()
        gr.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(0), Emu(0), Emu(100 * EMU), Emu(50 * EMU))
        elem, _ = _dispatch_shape(gr)
        assert elem["type"] == "group"
        assert [e["type"] for e in elem["elements"]] == ["shape"]

    def test_picture_increments_counter_from_offset(self, tmp_path):
        _, slide = _blank_slide()
        pic = slide.shapes.add_picture(
            _png_stream(), Emu(0), Emu(0), Emu(100 * EMU), Emu(50 * EMU))
        elem, counter = _dispatch_shape(pic, output_dir=tmp_path, slide_idx=0, img_counter=5)
        assert elem["src"] == "images/slide1_image6.png"
        assert counter == 6

    def test_media_routes_to_video(self, tmp_path):
        elem, counter = _dispatch_shape(_FakeVideoShape(), output_dir=tmp_path, slide_idx=0, img_counter=0)
        assert elem["type"] == "video"
        assert counter == 1

    def test_table_routes_to_table(self):
        _, slide = _blank_slide()
        table_shape = slide.shapes.add_table(
            2, 2, Emu(0), Emu(0), Emu(400 * EMU), Emu(100 * EMU))
        elem, _ = _dispatch_shape(table_shape)
        assert elem["type"] == "table"

    def test_chart_routes_to_chart(self):
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE
        _, slide = _blank_slide()
        cd = CategoryChartData()
        cd.categories = ["c1"]
        cd.add_series("s1", (1.0,))
        chart_shape = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, Emu(0), Emu(0), Emu(400 * EMU), Emu(300 * EMU), cd)
        elem, _ = _dispatch_shape(chart_shape)
        assert elem["type"] == "chart"

    def test_wordart_routes_to_raw_passthrough(self, tmp_path):
        _, slide = _blank_slide()
        tb = slide.shapes.add_textbox(Emu(10 * EMU), Emu(10 * EMU), Emu(100 * EMU), Emu(50 * EMU))
        tb.text_frame.text = "WordArt"
        body_pr = tb._element.find(f".//{{{_A}}}bodyPr")
        warp = etree.SubElement(body_pr, f"{{{_A}}}prstTxWarp")
        warp.set("prst", "textArchUp")
        elem, counter = _dispatch_shape(tb, output_dir=tmp_path, slide_idx=0, img_counter=0)
        assert counter == 0
        shape_xml = elem.pop("_shapeXml")
        assert shape_xml.startswith("<p:sp")
        assert "prstTxWarp" in shape_xml
        assert elem == {"type": "rawShape", "x": 10, "y": 10, "width": 100, "height": 50}

    def test_blipfill_shape_routes_to_image(self, tmp_path):
        _, slide = _blank_slide()
        sh = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Emu(400 * EMU), Emu(100 * EMU), Emu(150 * EMU), Emu(100 * EMU))
        _, rid = slide.part.get_or_add_image_part(_png_stream(size=(8, 8)))
        sp_pr = sh._element.spPr
        bf = etree.SubElement(sp_pr, f"{{{_A}}}blipFill")
        bl = etree.SubElement(bf, f"{{{_A}}}blip")
        bl.set(f"{{{_R}}}embed", rid)
        elem, counter = _dispatch_shape(sh, output_dir=tmp_path, slide_idx=0, img_counter=1)
        assert counter == 2
        assert elem == {
            "type": "image", "x": 400, "y": 100, "width": 150, "height": 100,
            "src": "images/slide1_image2.png", "fit": "cover",
        }


# ---------------------------------------------------------------------------
# Pipeline contract (deterministic corpus, exact JSON)
# ---------------------------------------------------------------------------


def _build_corpus_deck(path):
    """Standard 16:9 deck exercising shape/textbox/line/image/group/table/chart."""
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    prs, slide = _blank_slide()
    sh = slide.shapes.add_shape(
        MSO_SHAPE.CHEVRON, Emu(100 * EMU), Emu(100 * EMU), Emu(200 * EMU), Emu(100 * EMU))
    sh.text_frame.text = "Chev"
    tb = slide.shapes.add_textbox(Emu(400 * EMU), Emu(100 * EMU), Emu(300 * EMU), Emu(60 * EMU))
    tb.text_frame.text = "Corpus"
    slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Emu(100 * EMU), Emu(300 * EMU), Emu(300 * EMU), Emu(400 * EMU))
    slide.shapes.add_picture(
        _png_stream(size=(16, 12), color=(5, 5, 5)),
        Emu(800 * EMU), Emu(100 * EMU), Emu(160 * EMU), Emu(120 * EMU))
    gr = slide.shapes.add_group_shape()
    gr.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Emu(1200 * EMU), Emu(600 * EMU), Emu(100 * EMU), Emu(50 * EMU))
    gr.shapes.add_shape(
        MSO_SHAPE.OVAL, Emu(1350 * EMU), Emu(600 * EMU), Emu(60 * EMU), Emu(60 * EMU))
    tbl = slide.shapes.add_table(
        2, 2, Emu(100 * EMU), Emu(500 * EMU), Emu(400 * EMU), Emu(120 * EMU)).table
    tbl.cell(0, 0).text = "A"
    tbl.cell(0, 1).text = "B"
    tbl.cell(1, 0).text = "1"
    tbl.cell(1, 1).text = "2"
    cd = CategoryChartData()
    cd.categories = ["c1", "c2"]
    cd.add_series("s1", (1.0, 2.0))
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Emu(900 * EMU), Emu(400 * EMU), Emu(500 * EMU), Emu(300 * EMU), cd)
    prs.save(path)


_EXPECTED_CORPUS_ELEMENTS = [
    {"type": "shape", "shape": "chevron", "x": 100, "y": 100, "width": 200, "height": 100,
     "fill": "#4F81BD", "line": "#4F81BD", "lineWidth": 0.8, "text": "Chev",
     "fontSize": 18, "fontColor": "#FFFFFF", "align": "left", "verticalAlign": "middle"},
    {"type": "textbox", "x": 400, "y": 100, "width": 300, "height": 60,
     "fill": "none", "line": "none", "text": "Corpus", "marginLeft": 14,
     "marginRight": 14, "autoWidth": True, "_spAutoFit": True},
    {"type": "line", "preset": "line", "x1": 100, "y1": 300, "x2": 300, "y2": 400,
     "color": "none"},
    {"type": "image", "x": 800, "y": 100, "width": 160, "height": 120,
     "src": "images/slide1_image1.png", "link": None, "fit": "stretch"},
    {"type": "group", "x": 1200, "y": 600, "width": 210, "height": 60, "elements": [
        {"type": "shape", "x": 1200, "y": 600, "width": 100, "height": 50,
         "shape": "rectangle", "fill": "#4F81BD", "line": "#4F81BD", "lineWidth": 0.8},
        {"type": "shape", "x": 1350, "y": 600, "width": 60, "height": 60,
         "shape": "circle", "fill": "#4F81BD", "line": "#4F81BD", "lineWidth": 0.8},
    ]},
    {"type": "table", "x": 100, "y": 500, "width": 400, "height": 120,
     "colWidths": [200, 200], "rowHeights": [60, 60],
     "headers": [
         {"vertical-align": "top", "text": "{{#000000:A}}"},
         {"vertical-align": "top", "text": "{{#000000:B}}"},
     ],
     "rows": [[
         {"vertical-align": "top", "text": "{{#000000:1}}"},
         {"vertical-align": "top", "text": "{{#000000:2}}"},
     ]]},
    {"type": "chart", "x": 900, "y": 400, "width": 500, "height": 300,
     "chartType": "bar", "categories": ["c1", "c2"],
     "series": [{"name": "s1", "values": [1.0, 2.0]}],
     "categoryAxis": {"gridlines": False}, "legend": False},
]


class TestPipelineContract:
    def test_standard_16x9_corpus(self, tmp_path):
        from sdpm.engine.converter.pipeline import pptx_to_json

        pptx_path = tmp_path / "corpus.pptx"
        _build_corpus_deck(pptx_path)
        out_dir = tmp_path / "out"
        result = pptx_to_json(pptx_path, out_dir)

        assert result["fonts"] == {"halfwidth": "Calibri", "fullwidth": ""}
        assert result["defaultTextColor"] == "#000000"
        assert len(result["slides"]) == 1
        slide = result["slides"][0]
        assert slide["layout"] == "Blank"
        assert slide["masterIndex"] == 0

        elements = [dict(e) for e in slide["elements"]]
        # _chartXml embeds python-pptx-generated random axIds — assert shape, not value
        chart_xml = elements[6].pop("_chartXml")
        assert "<c:barChart>" in chart_xml
        assert elements == _EXPECTED_CORPUS_ELEMENTS

        # On-disk deck structure
        deck_meta = json.loads((out_dir / "deck.json").read_text())
        assert deck_meta == {
            "fonts": {"halfwidth": "Calibri", "fullwidth": ""},
            "defaultTextColor": "#000000",
            "autoSpacing": False,
            "slideSize": {"width": 1920, "height": 1080, "ptPerPx": 0.5},
        }
        assert (out_dir / "slides" / "slide-001.json").exists()
        assert sorted(p.name for p in (out_dir / "images").iterdir()) == ["slide1_image1.png"]

    def test_non_standard_4x3_normalizes_to_1920_basis(self, tmp_path):
        """A 4:3 deck (9144000 EMU wide) must still map slide-width → 1920px."""
        from sdpm.engine.converter.pipeline import pptx_to_json

        prs, slide = _blank_slide(width_emu=9144000, height_emu=6858000)
        emu43 = 9144000 / 1920
        tb = slide.shapes.add_textbox(
            Emu(round(480 * emu43)), Emu(round(270 * emu43)),
            Emu(round(960 * emu43)), Emu(round(100 * emu43)))
        tb.text_frame.text = "centered43"
        pptx_path = tmp_path / "deck43.pptx"
        prs.save(pptx_path)

        result = pptx_to_json(pptx_path, tmp_path / "out43")
        assert result["slides"][0]["elements"] == [{
            "type": "textbox", "x": 480, "y": 270, "width": 960, "height": 100,
            "fill": "none", "line": "none", "text": "centered43",
            "marginLeft": 19, "marginRight": 19, "autoWidth": True, "_spAutoFit": True,
        }]


# ---------------------------------------------------------------------------
# Import-path & signature compatibility
# ---------------------------------------------------------------------------

_PUBLIC_EXTRACTORS = {
    "extract_line_element": ["shape", "theme_colors", "color_mapping", "theme_styles"],
    "extract_freeform_element": ["shape", "theme_colors", "color_mapping", "builder_text_color"],
    "extract_shape_element": ["shape", "theme_colors", "color_mapping", "theme_styles", "builder_text_color"],
    "extract_textbox_element": ["shape", "theme_colors", "color_mapping", "theme_styles",
                                "is_placeholder", "builder_text_color"],
    "extract_video_element": ["shape", "output_dir", "slide_idx", "img_idx"],
    "extract_picture_element": ["shape", "output_dir", "slide_idx", "img_idx",
                                "theme_colors", "color_mapping"],
    "extract_group_element": ["shape", "theme_colors", "color_mapping", "theme_styles",
                              "output_dir", "slide_idx", "img_counter", "builder_text_color"],
    "_dispatch_shape": ["shape", "theme_colors", "color_mapping", "theme_styles",
                        "output_dir", "slide_idx", "img_counter", "builder_text_color", "pptx_path"],
}

# The converter facade re-exports these six (video is deliberately not public there)
_FACADE_EXTRACTORS = [
    "extract_shape_element", "extract_textbox_element", "extract_line_element",
    "extract_freeform_element", "extract_picture_element", "extract_group_element",
]


class TestImportCompatibility:
    @pytest.mark.parametrize("name,params", sorted(_PUBLIC_EXTRACTORS.items()))
    def test_elements_module_exposes_function_with_signature(self, name, params):
        import sdpm.engine.converter.elements as elements

        fn = getattr(elements, name)
        assert list(inspect.signature(fn).parameters) == params

    @pytest.mark.parametrize("name", _FACADE_EXTRACTORS)
    def test_converter_facade_reexports_same_objects(self, name):
        import sdpm.engine.converter as converter
        import sdpm.engine.converter.elements as elements

        assert getattr(converter, name) is getattr(elements, name)

    def test_slide_module_import_surface(self):
        """converter.slide imports these directly — they must stay importable."""
        from sdpm.engine.converter.elements import (  # noqa: F401
            _dispatch_shape,
            extract_picture_element,
            extract_textbox_element,
        )
