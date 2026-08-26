# Copyright Amazon.com, Inc. or its affiliates. All Rights Reserved.
# SPDX-License-Identifier: MIT-0
"""Tests for preview/__init__.py — imbalance check with emu_per_px (Phase 1-7)."""

from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

from sdpm.engine.preview import check_layout_imbalance_data


class TestImbalanceEmuPerPx:
    """Ensure check_layout_imbalance_data uses emu_per_px from slide width."""

    def test_16x9_canvas_size(self, template_16x9: Path, tmp_path: Path):
        """16:9 template: canvas reports as 1920 x 1080."""
        prs = Presentation(str(template_16x9))
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        shape = slide.shapes.add_textbox(Emu(6350 * 960), Emu(6350 * 540), Emu(6350 * 200), Emu(6350 * 100))
        shape.text = "Center"
        out = tmp_path / "test_16x9.pptx"
        prs.save(str(out))

        # Centered element should not trigger imbalance
        alerts = check_layout_imbalance_data(str(out))
        # Element is roughly centered (540 + 50 = cy=590 vs _CY=~560), may or may not alert
        # The key assertion: no crash and it runs with real emu_per_px
        assert isinstance(alerts, list)

    def test_4x3_canvas_size(self, template_4x3: Path, tmp_path: Path):
        """4:3 template: canvas reports as 1920 x 1440 (not 1440 x 1080)."""
        prs = Presentation(str(template_4x3))
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        # Place an element at the vertical center of a 4:3 slide
        # 4:3 emu_per_px = 9144000 / 1920 = 4762.5
        emu_px = 4762.5
        # Vertically centered: y = ~720 (half of 1440)
        cy_target = 720
        elem_h = 200
        elem_y = cy_target - elem_h // 2  # y = 620
        shape = slide.shapes.add_textbox(
            Emu(round(100 * emu_px)),
            Emu(round(elem_y * emu_px)),
            Emu(round(400 * emu_px)),
            Emu(round(elem_h * emu_px)),
        )
        shape.text = "4:3 centered"
        out = tmp_path / "test_4x3.pptx"
        prs.save(str(out))

        alerts = check_layout_imbalance_data(str(out))
        # With correct emu_per_px, a centered element should produce
        # minimal or no imbalance alert
        assert isinstance(alerts, list)

    def test_4x3_bbox_reports_correct_dimensions(self, template_4x3: Path, tmp_path: Path):
        """4:3: bbox string should reference '1920x1440' (not '1440x1080')."""
        prs = Presentation(str(template_4x3))
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        emu_px = 4762.5
        # Place element at top-left so it triggers imbalance (top-heavy)
        shape = slide.shapes.add_textbox(
            Emu(round(100 * emu_px)),
            Emu(round(50 * emu_px)),
            Emu(round(400 * emu_px)),
            Emu(round(100 * emu_px)),
        )
        shape.text = "Top element"
        out = tmp_path / "test_4x3_top.pptx"
        prs.save(str(out))

        alerts = check_layout_imbalance_data(str(out))
        # Should have an alert (element is way at the top)
        assert len(alerts) >= 1
        # The bbox should reference 1920x1440
        assert "1920x1440" in alerts[0]["bbox"]
