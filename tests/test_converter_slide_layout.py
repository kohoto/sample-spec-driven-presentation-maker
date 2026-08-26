# Copyright Amazon.com, Inc. or its affiliates. All Rights Reserved.
# SPDX-License-Identifier: MIT-0
"""Tests for converter/slide.py — detect_layout with emu_per_px (Phase 1-7)."""

from pathlib import Path

from pptx import Presentation

from sdpm.engine.converter.constants import conversion_scale
from sdpm.engine.converter.slide import detect_layout


class TestDetectLayoutEmuPerPx:
    """Ensure detect_layout uses get_emu_per_px() for left-accent-line detection."""

    def _make_slide_with_accent_line(self, prs, left_px, width_px, height_px):
        """Add a slide using 'Title Only' layout with a thin left accent line.

        The accent line check only fires when has_title=True and has_content=False,
        which is the 'Title Only' layout (index 5).
        """
        layout = prs.slide_layouts[5]  # "Title Only" — has TITLE but no BODY/OBJECT
        slide = prs.slides.add_slide(layout)
        from pptx.util import Emu as EmuUtil
        emu_px = prs.slide_width / 1920
        slide.shapes.add_shape(
            1,  # rectangle
            EmuUtil(round(left_px * emu_px)),
            EmuUtil(0),
            EmuUtil(round(width_px * emu_px)),
            EmuUtil(round(height_px * emu_px)),
        )
        return slide

    def test_16x9_detects_accent_line(self, template_16x9: Path):
        """16:9: left accent line (< 50px wide, > 500px tall) detected as content."""
        prs = Presentation(str(template_16x9))
        slide = self._make_slide_with_accent_line(prs, left_px=10, width_px=5, height_px=800)

        # Within default conversion scope (16:9), detect should find the accent line
        result = detect_layout(slide)
        assert result == "content"

    def test_4x3_detects_accent_line(self, template_4x3: Path):
        """4:3: left accent line detected correctly with proper emu_per_px scope."""
        prs = Presentation(str(template_4x3))
        slide = self._make_slide_with_accent_line(prs, left_px=10, width_px=5, height_px=800)

        # Must be within 4:3 conversion scope for correct threshold
        with conversion_scale(prs.slide_width):
            result = detect_layout(slide)
        assert result == "content"

    def test_4x3_accent_line_without_scope_uses_default(self, template_4x3: Path):
        """4:3 without scope: default emu_per_px (6350) is more permissive on left threshold."""
        prs = Presentation(str(template_4x3))
        # Place accent line at 10px in 4:3 coords (emu = 10 * 4762.5 = 47625).
        # Default scope threshold: 50 * 6350 = 317500. 47625 < 317500 passes.
        # Width threshold: 20 * 6350 = 127000. shape width = 5 * 4762.5 = 23812. passes.
        # Height threshold: 500 * 6350 = 3175000. shape height = 800 * 4762.5 = 3810000. passes.
        slide = self._make_slide_with_accent_line(prs, left_px=10, width_px=5, height_px=800)
        # In default scope, the EMU values from 4:3 are smaller for width/left
        # but the thresholds (with larger emu_per_px) are also larger — still passes.
        result = detect_layout(slide)
        assert result == "content"
