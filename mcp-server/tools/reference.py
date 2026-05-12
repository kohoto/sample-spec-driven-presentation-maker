# Copyright Amazon.com, Inc. or its affiliates. All Rights Reserved.
# SPDX-License-Identifier: MIT-0
"""Reference document access via Storage ABC — 3 categories × list/read + example search."""

import re
import tempfile
import time
from pathlib import Path
from typing import Any

from storage import Storage

# In-memory cache: key → (content, timestamp)
_cache: dict[str, tuple[Any, float]] = {}
CACHE_TTL: int = 3600  # 1 hour


def _parse_frontmatter(content: str) -> dict[str, str]:
    """Extract YAML frontmatter from markdown content.

    Args:
        content: Markdown text potentially starting with --- frontmatter ---.

    Returns:
        Dict of frontmatter key-value pairs.
    """
    match = re.match(r"^---\s*\n(.*?)\n---", content, re.DOTALL)
    if not match:
        return {}
    fm: dict[str, str] = {}
    for line in match.group(1).split("\n"):
        if ":" in line:
            key, _, val = line.partition(":")
            fm[key.strip()] = val.strip().strip('"').strip("'")
    return fm


def _list_category(storage: Storage, category: str) -> list[dict[str, str]]:
    """List all .md and .pptx files in a reference category with metadata.

    Args:
        storage: Storage backend instance.
        category: Category name (examples, workflows, guides).

    Returns:
        List of dicts with name, description, category fields.
    """
    cache_key = f"list:{category}"
    cached = _cache.get(cache_key)
    if cached and (time.time() - cached[1]) < CACHE_TTL:
        return cached[0]

    prefix = f"references/{category}/"
    files = storage.list_files(prefix=prefix)
    items: list[dict[str, str]] = []

    for f in files:
        if f.endswith(".md"):
            name = f.removeprefix(prefix).removesuffix(".md")
            try:
                content = storage.download_file(key=f).decode("utf-8")
                fm = _parse_frontmatter(content)
            except Exception:
                fm = {}
            items.append({
                "name": name,
                "description": fm.get("description", ""),
                "category": fm.get("category", ""),
            })
        elif f.endswith(".pptx"):
            name = f.removeprefix(prefix).removesuffix(".pptx")
            items.append({
                "name": name,
                "description": "(pptx — use read_examples with page specifier)",
                "category": "",
            })

    _cache[cache_key] = (items, time.time())
    return items


def _read_docs(storage: Storage, category: str, names: list[str]) -> list[dict[str, str]]:
    """Read one or more documents from a reference category.

    Supports .md (returned as text) and .pptx (rendered via get_pptx_notes).
    Names can include page specifiers for pptx: "common/patterns/3" or "common/patterns/all".

    Args:
        storage: Storage backend instance.
        category: Category name (examples, workflows, guides).
        names: List of document names. Supports "name", "name/page", "name/all".

    Returns:
        List of dicts with name and content.

    Raises:
        FileNotFoundError: If a requested document is not found.
    """
    results: list[dict[str, str]] = []
    for name in names:
        # Parse page specifier (e.g. "common/patterns/3" → file="common/patterns", pages=[3])
        pages = None
        has_page_specifier = False
        parts = name.rsplit("/", 1)
        file_name = name
        if len(parts) == 2 and (parts[1].isdigit() or parts[1] == "all"):
            file_name = parts[0]
            has_page_specifier = True
            if parts[1] == "all":
                pages = None
            else:
                page_num = int(parts[1])
                if page_num < 1:
                    raise ValueError(f"Invalid page number: {parts[1]} (must be >= 1)")
                pages = [page_num]

        # Try .md first, then .pptx
        md_key = f"references/{category}/{file_name}.md"
        pptx_key = f"references/{category}/{file_name}.pptx"

        md_content = None
        try:
            md_content = storage.download_file(key=md_key).decode("utf-8")
        except Exception:
            pass

        if md_content is not None:
            results.append({"name": file_name, "content": md_content})
            continue

        pptx_bytes = None
        try:
            pptx_bytes = storage.download_file(key=pptx_key)
        except Exception:
            pass

        if pptx_bytes is not None:
            if not has_page_specifier:
                # No page specifier — return slide description listing
                content = _list_pptx_descriptions_from_bytes(pptx_bytes)
            else:
                content = _render_pptx_from_bytes(pptx_bytes, pages=pages)
            results.append({"name": name, "content": content})
            continue

        # Not found — list available for error message
        available = [item["name"] for item in _list_category(storage, category)]
        raise FileNotFoundError(
            f"'{file_name}' not found in {category}/. Available: {', '.join(available)}"
        )
    return results


def _list_pptx_descriptions_from_bytes(pptx_bytes: bytes) -> str:
    """Return a listing of slide descriptions (1st line of speaker notes per slide).

    Args:
        pptx_bytes: Raw pptx file content.

    Returns:
        Text listing with page number and description per line.
    """
    from sdpm.reference import list_pptx_descriptions

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        tmp.write(pptx_bytes)
        tmp_path = Path(tmp.name)

    try:
        descriptions = list_pptx_descriptions(str(tmp_path))
        lines = [f"  {page:>3}  {desc}" for page, desc in descriptions]
        return "\n".join(lines)
    finally:
        tmp_path.unlink(missing_ok=True)


def _render_pptx_from_bytes(pptx_bytes: bytes, pages: list[int] | None = None) -> str:
    """Render pptx bytes to text via get_pptx_notes.

    Args:
        pptx_bytes: Raw pptx file content.
        pages: Optional list of 1-based page numbers. None = all pages.

    Returns:
        Rendered text output with page headers and notes.
    """
    from sdpm.reference import get_pptx_notes

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        tmp.write(pptx_bytes)
        tmp_path = Path(tmp.name)

    try:
        results = get_pptx_notes(tmp_path, pages=pages)
        lines: list[str] = []
        for page_num, notes in results:
            lines.append(f"## Page {page_num}\n\n{notes}\n")
        return "\n".join(lines)
    finally:
        tmp_path.unlink(missing_ok=True)


def list_styles(storage: Storage, user_id: str = "", include_all: bool = False) -> dict[str, Any]:
    """List available design styles with pin/source metadata.

    Combines builtin styles (references/examples/styles/) and user styles
    (user-styles/{user_id}/). Uses Engine filter_styles() for filtering.

    Args:
        storage: Storage backend instance.
        user_id: User ID for fetching user styles and pins. Empty = builtin only.
        include_all: If True, return all styles. If False, return pinned + user only
                     (falls back to all if no pins exist).

    Returns:
        Dict with styles list (name, description, pinned, source).
    """
    from sdpm.reference import filter_styles

    # 1. Builtin styles from resource bucket
    cache_key = "list:styles"
    cached = _cache.get(cache_key)
    if cached and (time.time() - cached[1]) < CACHE_TTL:
        builtin_styles = cached[0]
    else:
        prefix = "references/examples/styles/"
        files = storage.list_files(prefix=prefix)
        builtin_styles: list[dict[str, str]] = []
        for f in files:
            if not f.endswith(".html"):
                continue
            name = f.removeprefix(prefix).removesuffix(".html")
            description = ""
            try:
                content = storage.download_file(key=f).decode("utf-8")
                m = re.search(r"<title>(.*?)</title>", content, re.IGNORECASE)
                if m:
                    description = m.group(1).strip()
            except Exception:
                pass
            builtin_styles.append({"name": name, "description": description, "source": "builtin"})
        _cache[cache_key] = (builtin_styles, time.time())

    # 2. User styles from pptx bucket
    user_styles: list[dict[str, str]] = []
    if user_id:
        user_prefix = f"user-styles/{user_id}/"
        user_files = storage.list_files(prefix=user_prefix, bucket=storage.pptx_bucket)
        for f in user_files:
            if not f.endswith(".html"):
                continue
            name = f.removeprefix(user_prefix).removesuffix(".html")
            description = ""
            try:
                content = storage.download_file_from_pptx_bucket(key=f).decode("utf-8")
                m = re.search(r"<title>(.*?)</title>", content, re.IGNORECASE)
                if m:
                    description = m.group(1).strip()
            except Exception:
                pass
            user_styles.append({"name": name, "description": description, "source": "user"})

    # 3. Get pins
    pinned_names: list[str] = []
    if user_id:
        pinned_names = storage.get_style_pins(user_id)

    # 4. Filter via Engine
    all_styles = user_styles + builtin_styles
    filtered = filter_styles(all_styles, pinned_names, include_all)
    return {"styles": filtered}


def read_examples(names: list[str], storage: Storage) -> dict[str, Any]:
    """Read one or more example documents.

    Without page specifier (e.g. "patterns"), returns a listing of slide
    descriptions. With page specifier (e.g. "patterns/3"), returns full content.

    Args:
        names: Example names (e.g. ["patterns", "patterns/3", "components/all"]).
        storage: Storage backend instance.

    Returns:
        Dict with documents list.
    """
    return {"documents": _read_docs(storage, "examples", names)}


def list_workflows(storage: Storage) -> dict[str, Any]:
    """List all workflow documents.

    Args:
        storage: Storage backend instance.

    Returns:
        Dict with items list.
    """
    return {"items": _list_category(storage, "workflows")}


def read_workflows(names: list[str], storage: Storage) -> dict[str, Any]:
    """Read one or more workflow documents.

    Args:
        names: Workflow names (e.g. ["create-new-2-build"]).
        storage: Storage backend instance.

    Returns:
        Dict with documents list.
    """
    return {"documents": _read_docs(storage, "workflows", names)}


def list_guides(storage: Storage) -> dict[str, Any]:
    """List all guide documents.

    Args:
        storage: Storage backend instance.

    Returns:
        Dict with items list.
    """
    return {"items": _list_category(storage, "guides")}


def read_guides(names: list[str], storage: Storage) -> dict[str, Any]:
    """Read one or more guide documents.

    Args:
        names: Guide names (e.g. ["design-rules"]).
        storage: Storage backend instance.

    Returns:
        Dict with documents list.
    """
    return {"documents": _read_docs(storage, "guides", names)}


def search_examples(query: str, storage: Storage, limit: int = 5) -> dict[str, Any]:
    """Search pptx example slides by keywords in speaker notes.

    Downloads pptx files from S3 to a temp directory and delegates to
    the engine's search_examples function.

    Args:
        query: Space-separated keywords.
        storage: Storage backend instance.
        limit: Maximum results to return.

    Returns:
        Dict with results list.
    """
    # Download pptx examples to temp dir, reusing cache
    cache_key = "examples_tmpdir"
    cached = _cache.get(cache_key)
    if cached and (time.time() - cached[1]) < CACHE_TTL:
        tmp_examples_dir = Path(cached[0])
    else:
        tmp_examples_dir = Path(tempfile.mkdtemp(prefix="pptx-examples-"))
        prefix = "references/examples/"
        files = storage.list_files(prefix=prefix)
        for f in files:
            if not f.endswith(".pptx"):
                continue
            rel = f.removeprefix(prefix)
            local_path = tmp_examples_dir / rel
            local_path.parent.mkdir(parents=True, exist_ok=True)
            data = storage.download_file(key=f)
            local_path.write_bytes(data)
        _cache[cache_key] = (str(tmp_examples_dir), time.time())

    # Monkey-patch the examples dir for search_examples (it uses __file__ relative path)
    # Instead, call the underlying logic directly
    from pptx import Presentation as PptxPresentation

    queries = query.lower().split()
    raw_results: list[tuple[int, str, int, str]] = []

    for pptx_file in sorted(tmp_examples_dir.rglob("*.pptx")):
        if "must-read" in pptx_file.stem:
            continue
        cat = pptx_file.parent.name
        stem = pptx_file.stem
        try:
            prs = PptxPresentation(str(pptx_file))
        except Exception:
            continue
        for si, slide in enumerate(prs.slides):
            if not slide.has_notes_slide:
                continue
            notes = slide.notes_slide.notes_text_frame.text.replace('\x0B', '\n')
            if not notes.strip():
                continue
            notes_lower = notes.lower()
            match_count = sum(1 for q in queries if q in notes_lower)
            if match_count == 0:
                continue
            desc = ""
            for line in notes.splitlines():
                if line.strip():
                    desc = line.strip()
                    break
            raw_results.append((match_count, f"{cat}/{stem}", si + 1, desc))

    raw_results.sort(key=lambda x: (-x[0], x[1]))
    if not raw_results:
        return {"results": []}
    limited = raw_results[:limit]
    min_score = limited[-1][0]
    limited = [r for r in raw_results if r[0] >= min_score]
    return {
        "results": [
            {"path": r[1], "page": r[2], "description": r[3], "score": f"{r[0]}/{len(queries)}"}
            for r in limited
        ],
    }
